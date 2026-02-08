from __future__ import annotations

import os
import sys
import asyncio
import random
import re
import unicodedata
import threading
import queue
import logging
import time
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
from urllib.parse import urlparse, urlunparse, quote, parse_qs
import tempfile
import zipfile
import io

# =========================
# Dependencias PDF / Word / PPT
# =========================
try:
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover
    fitz = None  # type: ignore

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None  # type: ignore

# 🔹 NUEVO: PPTX
try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore

# =========================
# Dependencia para Descarga Masiva
# =========================
try:
    import requests
except ImportError:  # pragma: no cover
    requests = None  # type: ignore

import pandas as pd
import streamlit as st

try:
    import httpx
except ImportError:
    httpx = None  # type: ignore


# ======================================================
# LOGGING
# ======================================================

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.StreamHandler(sys.stdout),
    ],
)
logger = logging.getLogger(__name__)


# ======================================================
# CONFIG / CONSTANTES
# ======================================================

APP_TITLE = "UTP - Broken Link Checker"
APP_ICON = "🔗"

MODULES = [
    "Home",
    "Report Broken Link",  # Solo 2 módulos
]



DEFAULT_TIMEOUT_S = 15.0
DEFAULT_CONCURRENCY_GLOBAL = 30
DEFAULT_CONCURRENCY_PER_HOST = 6
DEFAULT_RETRIES = 3

DEFAULT_MAX_BYTES = 200_000
DEFAULT_RANGE_BYTES = 131_072

# --------- Descarga Masiva ----------
MAX_INTENTOS_DESCARGA = 7
CHUNK_SIZE = 1024 * 256  # 256 KB
REQUEST_HEADERS = {
    "User-Agent": "Mozilla/5.0 (compatible; UTP-FileDownloader/1.0)",
    "Accept": "*/*",
    "Accept-Encoding": "identity",
}
DESC_EXT_PERMITIDAS = (".ppt", ".pptx", ".pdf", ".doc", ".docx")

# ======================================================
# PATRONES DE SOFT-404 (incluye caso El Peruano)
# ======================================================

SOFT_404_PATTERNS = [
    # Básicos
    r"\b404\b",
    # r"not\s+found",  # eliminado por ser demasiado genérico
    r"page\s+not\s+found",
    r"file\s+not\s+found",
    r"document\s+not\s+found",
    r"resource\s+not\s+found",

    # Español
    r"p[aá]gina\s+no\s+encontrada",
    r"la\s+p[aá]gina\s+no\s+existe",
    r"no\s+se\s+encontr[oó]",
    r"no\s+podemos\s+encontrar",

    # Genéricos (pero más específicos que 'not found')
    r"error\s*404",
    r"http\s+404",
    r"error\s*:\s*page\s+not\s+available",
    r"the\s+page\s+you\s+requested\s+was\s+not\s+found",
    r"sorry,\s+this\s+page\s+doesn't\s+exist",
    r"oops!?\s*.*not\s+found",

    # Casos específicos como El Peruano
    r"the\s+specified\s+url\s+cannot\s+be\s+found",
    r"url\s+cannot\s+be\s+found",
    r"cannot\s+be\s+found",
    r"recurso\s+no\s+disponible",
    r"contenido\s+no\s+disponible",

    # YouTube
    r"video\s+unavailable",
    r"this\s+video\s+is\s+unavailable",
    r"this\s+video\s+isn't\s+available",
    r"video\s+no\s+disponible",

    # Repositorios / sistemas documentales
    r"item\s+not\s+found",
    r"handle\s+not\s+found",
    r"bitstream\s+not\s+found",
]
SOFT_404_RE = re.compile("|".join(SOFT_404_PATTERNS), re.IGNORECASE)

# Extensiones binarias candidatas
BINARY_EXTS = {
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".zip", ".rar", ".7z",
    ".ppt", ".pptx", ".mp4", ".mp3", ".avi", ".mov", ".wmv", ".mkv",
    ".png", ".jpg", ".jpeg", ".gif", ".tif", ".tiff", ".bmp", ".svg",
    ".webp", ".ico", ".exe", ".dmg", ".apk",
}

# 🔴 Regla especial: cualquier enlace que contenga 'canvas.utp'
# debe considerarse SIEMPRE como ROTO en el reporte.
CANVAS_UTP_KEYWORD = "canvas.utp"

# Mapeo de extensiones a Content-Type esperado
EXPECTED_CONTENT_TYPES = {
    ".pdf": ["application/pdf"],
    ".doc": ["application/msword"],
    ".docx": ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"],
    ".xls": ["application/vnd.ms-excel"],
    ".xlsx": ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"],
    ".zip": ["application/zip", "application/x-zip-compressed"],
    ".rar": ["application/x-rar-compressed"],
    ".ppt": ["application/vnd.ms-powerpoint"],
    ".pptx": ["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
    ".png": ["image/png"],
    ".jpg": ["image/jpeg"],
    ".jpeg": ["image/jpeg"],
    ".gif": ["image/gif"],
    ".mp4": ["video/mp4"],
    ".mp3": ["audio/mpeg"],
}

# ======================================================
# MEJORAS V5: USER-AGENTS Y CONFIGURACIONES POR DOMINIO
# ======================================================

# User-Agents realistas para diferentes navegadores
USER_AGENTS = [
    # Chrome (Windows)
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    # Firefox (Windows)
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:123.0) "
    "Gecko/20100101 Firefox/123.0",
    # Safari (macOS)
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 14_3) AppleWebKit/605.1.15 "
    "(KHTML, like Gecko) Version/17.2 Safari/605.1.15",
    # Chrome (macOS)
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    # Edge (Windows)
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36 Edg/122.0.0.0",
]

# Lista blanca de dominios confiables (menos agresivo con soft-404)
TRUSTED_DOMAINS = [
    "canva.com",
    "crehana.com",
    "openstax.org",
    "coca-cola.com",
    "wikipedia.org",
    "youtube.com",
    "youtu.be",
    "github.com",
    "stackoverflow.com",
    "google.com",
    "microsoft.com",
    "amazon.com",
    "bbc.com",
]


DOMAIN_CONFIGS: Dict[str, Dict[str, Any]] = {
    "facebook.com": {
        "requires_cookies": True,
        "accept_codes": [200, 302, 400],
        "additional_headers": {
            "Accept": (
                "text/html,application/xhtml+xml,application/xml;q=0.9,"
                "image/avif,image/webp,*/*;q=0.8"
            ),
            "Accept-Language": "es-PE,es;q=0.9,en-US;q=0.8,en;q=0.7",
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "none",
            "Sec-Fetch-User": "?1",
            "Upgrade-Insecure-Requests": "1",
        },
    },
    "linkedin.com": {
        "requires_cookies": True,
        # 999 es típico de LinkedIn cuando bloquea bots, pero el recurso existe
        "accept_codes": [200, 302, 999],
        "additional_headers": {
            "Accept": (
                "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            ),
            "Accept-Language": "es-PE,es;q=0.9,en;q=0.8",
        },
    },
    "canva.com": {
        # 403 aquí significa "necesitas login", pero el recurso existe
        "accept_codes": [200, 403],
        "additional_headers": {
            "Accept": (
                "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8"
            ),
            "Referer": "https://www.google.com/",
        },
    },
    "crehana.com": {
        "trusted_domain": True,
        "accept_codes": [200, 206],
    },
    "openstax.org": {
        "trusted_domain": True,
        "accept_codes": [200],
    },
    "coca-cola.com": {
        "trusted_domain": True,
        "accept_codes": [200, 206],
    },

    # Biblioteca UTP
    "tubiblioteca.utp.edu.pe": {
        "trusted_domain": True,
        # redirecciones a login + 200 son válidos
        "accept_codes": [200, 302, 303, 307, 308],
        # este dominio puede dar problemas de certificado → lo ignoramos
        "skip_ssl_verify": True,
    },

    # Wikipedia: a veces devuelve 403 a scripts, pero la página existe
    "wikipedia.org": {
        "trusted_domain": True,
        "accept_codes": [200, 301, 302, 403, 429],
    },

    # YouTube (vídeos, incluidos privados / eliminados)
    "youtube.com": {
        "trusted_domain": True,
        "accept_codes": [200, 301, 302, 303, 307, 308, 403, 429],
    },
    "youtu.be": {
        "trusted_domain": True,
        "accept_codes": [200, 301, 302, 303, 307, 308, 403, 429],
    },

    # X / Twitter: página puede devolver 200 aunque la cuenta o el post no exista
    "x.com": {
        "accept_codes": [200, 301, 302, 303, 307, 308],
    },
    "twitter.com": {
        "accept_codes": [200, 301, 302, 303, 307, 308],
    },

     # 🔴 NUEVO: Google Sites – redirecciones dentro del mismo sitio
    "sites.google.com": {
        "trusted_domain": False,
        "accept_codes": [200, 301, 302, 303, 307, 308],
    },

    # 🔴 NUEVO: Universidad de Granada – posibles problemas de SSL
    "ugr.es": {
        "trusted_domain": True,
        "accept_codes": [200, 301, 302, 303, 307, 308],
        "skip_ssl_verify": True,
    },

    # 🔴 NUEVO: WordPress / Cloudflare – 403/503 para bots, pero recurso existe
    "wordpress.com": {
        "trusted_domain": True,
        "accept_codes": [200, 301, 302, 303, 307, 308, 403, 503],
    },

}

# ======================================================
# PATRONES DE SOFT-404 MEJORADOS (V5)
# ======================================================

# ======================================================
# PATRONES DE SOFT-404 MEJORADOS (V5) - REFORZADO
# ======================================================

# Patrones muy específicos que indican con alta confianza que la página
# no existe o el recurso ya no está disponible (soft-404).
SOFT_404_STRONG_PATTERNS = [
    # Inglés
    r"\b404\s+error\b",
    r"\berror\s+404\b",
    r"page\s+not\s+found",
    r"file\s+not\s+found",
    r"document\s+not\s+found",
    r"resource\s+not\s+found",
    r"the\s+page\s+you.*(?:requested|looking\s+for).*not\s+found",
    r"sorry.*page.*doesn't\s+exist",

    # 🔴 NUEVO: muy típico de Apache / Google / otros
    r"the\s+requested\s+url\s+was\s+not\s+found\s+on\s+this\s+server",
    r"la\s+url\s+solicitada\s+no\s+se\s+ha\s+encontrado\s+en\s+este\s+servidor",
    r"no\s+se\s+ha\s+encontrado\s+la\s+url\s+solicitada",
    r"no\s+se\s+ha\s+encontrado\s+el\s+sitio\s+web",

    # Español
    r"p[aá]gina\s+no\s+encontrada",
    r"la\s+p[aá]gina\s+no\s+existe",
    r"la\s+p[aá]gina\s+que\s+buscas\s+no\s+existe",
    r"esta\s+p[aá]gina\s+no\s+existe",
    r"esta\s+p[aá]gina\s+no\s+est[aá]\s+disponible",
    r"recurso\s+no\s+disponible",
    r"contenido\s+no\s+disponible",
    r"the\s+specified\s+url\s+cannot\s+be\s+found",

    # Repositorios / sistemas documentales
    r"item\s+not\s+found",
    r"handle\s+not\s+found",
    r"bitstream\s+not\s+found",

    # YouTube (vídeo privado / eliminado)
    r"this\s+video\s+(?:isn't\s+available|is\s+unavailable)",
    r"video\s+no\s+disponible",
    r"este\s+video\s+es\s+privado",
    r"este\s+video\s+ha\s+sido\s+eliminado",

    # X / Twitter (página que no existe)
    r"esta\s+cuenta\s+no\s+existe",
    r"la\s+cuenta\s+no\s+existe",
    r"esta\s+p[aá]gina\s+no\s+existe\.\s*intenta\s+hacer\s+otra\s+b[uú]squeda",
]

SOFT_404_STRONG_RE = re.compile("|".join(SOFT_404_STRONG_PATTERNS), re.IGNORECASE)


# Patrones que indican contenido real (artículo, post, etc.)
VALID_CONTENT_PATTERNS = [
    r"<article",
    r"<main",
    r"class=['\"](?:content|article|post|entry)",
    r"<div\s+id=['\"](?:content|main|article)",
    r"og:type",
    r"twitter:card",
]
VALID_CONTENT_RE = re.compile("|".join(VALID_CONTENT_PATTERNS), re.IGNORECASE)

# ======================================================
# HELPERS V5 PARA DOMINIOS / SCORING
# ======================================================

def _get_domain_config(url: str) -> Dict[str, Any]:
    """Config específica para un dominio (facebook, canva, etc.)."""
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        for key in DOMAIN_CONFIGS:
            if key in domain:
                return DOMAIN_CONFIGS[key]
        return {}
    except Exception:
        return {}


def _is_trusted_domain(url: str) -> bool:
    """¿El dominio está en la lista blanca o marcado como 'trusted_domain'?"""
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()

        for trusted in TRUSTED_DOMAINS:
            if trusted in domain:
                return True

        config = _get_domain_config(url)
        return bool(config.get("trusted_domain", False))
    except Exception:
        return False


def _get_random_user_agent() -> str:
    """Devuelve un User-Agent aleatorio de la lista."""
    return random.choice(USER_AGENTS)


def _build_headers_for_domain(url: str) -> Dict[str, str]:
    """
    Headers "tipo navegador" + headers adicionales por dominio.
    Se usan como headers base del cliente httpx.
    """
    headers = {
        "user-agent": _get_random_user_agent(),
        "accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "accept-language": "es-PE,es;q=0.9,en;q=0.8",
        "accept-encoding": "gzip, deflate, br",
        "dnt": "1",
        "connection": "keep-alive",
        "upgrade-insecure-requests": "1",
    }

    config = _get_domain_config(url)
    if config.get("additional_headers"):
        headers.update(config["additional_headers"])

    return headers


def _is_valid_status_code(url: str, status_code: int) -> bool:
    """
    Determina si un código HTTP puede considerarse "válido" para un dominio.
    - Usa accept_codes específicos por dominio (facebook, canva, etc.).
    - Por defecto, cualquier código < 400 se considera válido.
    """
    config = _get_domain_config(url)
    accept_codes = config.get("accept_codes", [])

    if accept_codes and status_code in accept_codes:
        return True

    # Por defecto, 2xx y 3xx son válidos
    return status_code < 400


def _calculate_content_score(text: str, url: str) -> int:
    """
    Score basado en el contenido:
    - Score alto => página muy probablemente válida.
    - Score muy negativo => probable soft-404.
    """
    score = 0
    chunk = text[:5000]
    content_len = len(text)

    # Penalización por patrones fuertes de error 404
    if SOFT_404_RE.search(chunk):
        score -= 30

    # Bonificación por estructuras típicas de contenido real
    if VALID_CONTENT_RE.search(chunk):
        score += 20

    # Bonificación por longitud (las páginas de error suelen ser cortas)
    if content_len > 10000:
        score += 15
    elif content_len > 5000:
        score += 10
    elif content_len > 2000:
        score += 5

    # Dominio confiable => más bonificación
    if _is_trusted_domain(url):
        score += 25

    # Metadatos típicos de páginas reales
    head_chunk = text[:2000]
    if "og:title" in head_chunk or "twitter:title" in head_chunk:
        score += 10

    # Contenido demasiado corto
    if content_len < 500:
        score -= 10

    return score

# ======================================================
# VALIDADORES ESPECÍFICOS (LINK CHECKER)
# ======================================================

def validate_youtube_url(url: str) -> Tuple[bool, str]:
    if "youtube.com/watch" not in url and "youtu.be/" not in url:
        return True, ""

    try:
        parsed = urlparse(url)

        if "youtube.com" in parsed.netloc:
            query_params = parse_qs(parsed.query)
            if "v" not in query_params:
                return False, "URL de YouTube sin parámetro 'v'"
            video_id = query_params["v"][0]
        elif "youtu.be" in parsed.netloc:
            video_id = parsed.path.lstrip("/").split("/")[0]
        else:
            return True, ""

        if len(video_id) != 11:
            return False, f"ID de YouTube inválido (longitud {len(video_id)}, esperado 11)"
        if not re.match(r'^[A-Za-z0-9_-]{11}$', video_id):
            return False, "ID de YouTube contiene caracteres inválidos"

        return True, ""

    except Exception as e:
        return False, f"Error al validar YouTube: {str(e)}"

def validate_x_twitter_url(url: str) -> Tuple[bool, str]:
    """
    Valida URLs de X/Twitter.

    Consideramos sospechoso un enlace del estilo:
        https://x.com/<usuario>/status
    o
        https://twitter.com/<usuario>/status

    Es sintácticamente válido, pero no apunta a un tweet concreto
    porque le falta el ID numérico.
    """
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()

        if "x.com" in domain or "twitter.com" in domain:
            # Quitamos segmentos vacíos
            segments = [s for s in parsed.path.split("/") if s]

            # Esperamos algo como: /usuario/status/<id>
            if len(segments) >= 2 and segments[1].startswith("status"):
                # Falta el ID, o no es numérico → lo marcamos como inválido
                if len(segments) < 3 or not re.fullmatch(r"\d{5,}", segments[2]):
                    return False, "URL de X/Twitter sin ID de tweet (status incompleto)"

        return True, ""
    except Exception as e:
        return False, f"Error al validar X/Twitter: {str(e)}"

def validate_google_search_url(url: str) -> Tuple[bool, str]:
    """
    Marca como inválidas las URLs que apuntan a resultados de búsqueda de Google, por ejemplo:
        https://www.google.com/search?q=movimiento+parabolico&...
    Estas URLs son producto de una consulta y no son un recurso final estable.
    """
    try:
        parsed = urlparse(url)
        domain = parsed.netloc.lower()
        path = parsed.path or ""

        # Solo nos interesan las búsquedas tipo /search de Google
        if "google.com" in domain and path.startswith("/search"):
            # Si hay parámetro de búsqueda 'q', asumimos que es una consulta genérica
            qs = parse_qs(parsed.query)
            if "q" in qs and qs["q"] and any(qs["q"]):
                return False, (
                    "URL de búsqueda de Google (resultado de una consulta; "
                    "usar el enlace directo al recurso y no la búsqueda)"
                )

        return True, ""
    except Exception as e:
        return False, f"Error al validar búsqueda de Google: {str(e)}"



def validate_url_structure(url: str) -> Tuple[bool, str]:
    if "\\" in url:
        return False, "URL contiene backslash (\\) - carácter inválido"

    if " " in url and "%20" not in url:
        return False, "URL contiene espacios sin encodear"

    if not url.startswith(("http://", "https://")):
        return False, "URL debe comenzar con http:// o https://"

    try:
        parsed = urlparse(url)

        if not parsed.netloc:
            return False, "URL sin dominio válido"

        if "." not in parsed.netloc and "localhost" not in parsed.netloc.lower():
            return False, "Dominio inválido (falta extensión)"

        return True, ""

    except Exception as e:
        return False, f"Error de estructura: {str(e)}"


def validate_content_type_match(url: str, content_type: str) -> Tuple[bool, str]:
    if not content_type:
        return True, ""

    try:
        parsed = urlparse(url)
        path = parsed.path.lower()

        for ext, expected in EXPECTED_CONTENT_TYPES.items():
            if path.endswith(ext):
                ct_clean = content_type.lower().split(";")[0].strip()
                if ct_clean not in [e.lower() for e in expected]:
                    if "text/html" in ct_clean or "text/plain" in ct_clean:
                        return False, f"Archivo {ext} devuelve {ct_clean} (esperado {expected[0]})"
                break

        return True, ""

    except Exception:
        return True, ""


# ======================================================
# UI HELPERS
# ======================================================

def apply_global_styles():
    st.markdown(
        """
        <style>
        [data-testid="stAppViewContainer"] { background: #f3f4f6; }
        [data-testid="stSidebar"] { background: #f9fafb; }

        .utp-hero {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 18px;
            padding: 1.8rem 2.4rem;
            color: #ffffff;
            margin-bottom: 1.8rem;
            box-shadow: 0 18px 40px rgba(76, 81, 191, 0.35);
            display: flex;
            align-items: center;
            gap: 1.0rem;
        }
        .utp-hero-icon {
            width: 3.1rem;
            height: 3.1rem;
            border-radius: 999px;
            background: rgba(255,255,255,0.18);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 2.0rem;
        }
        .utp-hero-title { font-weight: 800; font-size: 1.8rem; margin-bottom: 0.15rem; }
        .utp-hero-sub { font-size: 0.92rem; opacity: 0.96; line-height: 1.4; }

        .utp-sidebar-brand {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 18px;
            padding: 1.0rem 1.1rem;
            color: #ffffff;
            box-shadow: 0 14px 32px rgba(76, 81, 191, 0.35);
            margin-bottom: 1.1rem;
        }
        .utp-sidebar-brand-title {
            font-weight: 800;
            font-size: 1.05rem;
            margin-bottom: 0.2rem;
            display: flex;
            align-items: center;
            gap: 0.4rem;
        }
        .utp-sidebar-brand-subtitle { font-size: 0.82rem; opacity: 0.92; }

        .utp-card {
            border-radius: 14px;
            border: 1px solid #e5e7eb;
            padding: 1.1rem 1.3rem 1.15rem 1.3rem;
            margin-bottom: 1.0rem;
            background: #ffffff;
            box-shadow: 0 10px 25px rgba(15,23,42,0.05);
        }

        .utp-step-row {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 0.7rem;
        }
        .utp-step-main {
            display: flex;
            align-items: center;
            gap: 0.55rem;
            font-size: 1.0rem;
            font-weight: 700;
            color: #111827;
        }
        .utp-step-number {
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 26px;
            height: 26px;
            border-radius: 999px;
            background: #4f46e5;
            color: #ffffff;
            font-size: 0.9rem;
            font-weight: 700;
            box-shadow: 0 3px 8px rgba(79,70,229,0.45);
        }
        .utp-step-status {
            padding: 0.18rem 0.7rem;
            border-radius: 999px;
            font-size: 0.78rem;
            font-weight: 600;
            border: 1px solid transparent;
            white-space: nowrap;
        }
        .utp-step-status-ok { background-color: #dcfce7; color: #166534; border-color: #bbf7d0; }
        .utp-step-status-warn { background-color: #ffedd5; color: #9a3412; border-color: #fed7aa; }
        .utp-step-status-error { background-color: #fee2e2; color: #b91c1c; border-color: #fecaca; }

        .utp-step-header-simple {
            display: flex;
            align-items: center;
            gap: 0.55rem;
            font-size: 1.0rem;
            font-weight: 700;
            color: #111827;
            margin-bottom: 0.7rem;
        }
        .utp-step-header-simple .utp-step-number {
            width: 26px;
            height: 26px;
            border-radius: 999px;
            background: #4f46e5;
            color: #ffffff;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            font-size: 0.9rem;
            font-weight: 700;
        }

        .stButton>button {
            border-radius: 999px;
            font-weight: 700;
            padding: 0.6rem 1.3rem;
            border: none;
            transition: all 0.2s ease;
        }
        .stButton>button:hover {
            transform: translateY(-1px);
            box-shadow: 0 10px 25px rgba(79,70,229,0.45);
        }

        .stDataFrame { border-radius: 10px; border: 1px solid #e5e7eb; }
        /* 🔁 Botón de reset justo debajo del hero, visualmente dentro de la franja */
        .hero-reset-anchor + div[data-testid="stButton"] {
            margin-top: -3.0rem;         /* lo sube encima del hero */
            margin-bottom: 0.6rem;
            display: flex;
            justify-content: flex-end;   /* alineado a la derecha */
            padding-right: 2.4rem;       /* mismo padding horizontal del hero */
        }

        /* Estilo del botón como icono circular rojo */
        .hero-reset-anchor + div[data-testid="stButton"] > button {
            position: relative;
            width: 44px;
            height: 44px;
            border-radius: 999px;
            padding: 0;
            background-color: #ff1654;   /* rojo parecido al icono adjunto */
            color: transparent;          /* oculta el texto 'Reiniciar' */
            box-shadow: 0 12px 28px rgba(15,23,42,0.45);
        }

        .hero-reset-anchor + div[data-testid="stButton"] > button::before {
            content: "↻";                /* icono de recarga */
            position: absolute;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            font-size: 22px;
            color: #ffffff;
        }

        .hero-reset-anchor + div[data-testid="stButton"] > button:hover {
            transform: translateY(-1px);
            box-shadow: 0 16px 32px rgba(15,23,42,0.55);
        }
        /* 🔵 Barra de progreso tipo "task" (progress-bar-ui-task) */
        .progress-bar-ui-task {
            width: 100%;
            height: 24px;
            border-radius: 999px;
            background: #e0f2fe; /* azul muy claro */
            box-shadow: inset 0 0 0 1px rgba(37,99,235,0.18);
            overflow: hidden;
            position: relative;
            margin: 0.35rem 0 0.15rem 0;
        }

        .progress-bar-ui-task__inner {
            position: relative;
            height: 100%;
            width: 0%;
            border-radius: inherit;
            background: linear-gradient(90deg, #0ea5e9 0%, #2563eb 60%, #1d4ed8 100%);
            display: flex;
            align-items: center;
            transition: width 0.25s ease-out;
        }

        .progress-bar-ui-task__label {
            padding-left: 0.9rem;
            font-size: 0.82rem;
            font-weight: 700;
            color: #ffffff;
            letter-spacing: 0.03em;
            white-space: nowrap;
            text-shadow: 0 1px 2px rgba(15,23,42,0.35);
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_sidebar_header():
    st.markdown(
        f"""
        <div class="utp-sidebar-brand">
            <div class="utp-sidebar-brand-title">
                <span>🔗</span><span>{APP_TITLE}</span>
            </div>
            <div class="utp-sidebar-brand-subtitle">
                Detección avanzada de links rotos con reporte único en Excel
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_hero(title: str, subtitle: str, icon: str = "🔗"):
    st.markdown(
        f"""
        <div class="utp-hero">
            <div class="utp-hero-icon">{icon}</div>
            <div>
                <div class="utp-hero-title">{title}</div>
                <div class="utp-hero-sub">{subtitle}</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_step_header_html(step_label: str, title: str, status: str) -> str:
    map_text = {"ok": "Listo", "warn": "Pendiente", "error": "Falta"}
    map_class = {"ok": "utp-step-status-ok", "warn": "utp-step-status-warn", "error": "utp-step-status-error"}
    status_text = map_text.get(status, "Pendiente")
    status_class = map_class.get(status, "utp-step-status-warn")
    return f"""
    <div class="utp-step-row">
        <div class="utp-step-main">
            <span class="utp-step-number">{step_label}</span>
            <span>{title}</span>
        </div>
        <div class="utp-step-status {status_class}">{status_text}</div>
    </div>
    """


def render_simple_step_header(step_label: str, title: str):
    st.markdown(
        f"""
        <div class="utp-step-header-simple">
            <span class="utp-step-number">{step_label}</span>
            <span>{title}</span>
        </div>
        """,
        unsafe_allow_html=True,
    )


def ui_card_open():
    st.markdown('<div class="utp-card">', unsafe_allow_html=True)


def ui_card_close():
    st.markdown("</div>", unsafe_allow_html=True)

def render_progress_bar_ui_task(placeholder, progress: float, label: Optional[str] = None):
    """
    Renderiza una barra tipo 'progress-bar-ui-task' dentro de un placeholder de Streamlit.

    - `progress` puede venir en rango 0–1 (0.0 = 0%, 1.0 = 100%)
      o 0–100 (60 = 60%).
    """
    if placeholder is None:
        return

    try:
        value = float(progress)
    except (TypeError, ValueError):
        value = 0.0

    # Aceptar tanto 0–1 como 0–100
    if value > 1.0:
        if value <= 100.0:
            value = value / 100.0
        else:
            value = 1.0
    value = max(0.0, min(1.0, value))

    percent_str = f"{value * 100:.1f}%"
    label_text = label or percent_str

    placeholder.markdown(
        f"""
        <div class="progress-bar-ui-task">
            <div class="progress-bar-ui-task__inner" style="width: {value * 100:.4f}%;">
                <span class="progress-bar-ui-task__label">{label_text}</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def _style_status_dataframe(df: pd.DataFrame):
    """
    Devuelve un Styler con la columna Status coloreada:
    - ACTIVO: verde
    - ROTO: rojo
    """
    if "Status" not in df.columns:
        return df.style

    def _color_status(val):
        v = str(val).upper()
        if v == "ACTIVO":
            return "background-color: #bbf7d0;"  # verde suave
        if v == "ROTO":
            return "background-color: #fecaca;"  # rojo suave
        return ""

    return df.style.applymap(_color_status, subset=["Status"])

def _render_status_summary(df_out: pd.DataFrame):
    """
    Muestra:
    - Métricas de ACTIVOS y ROTOS
    - Tabla principal (Nombre del Archivo | Link | Status)
    - Expander con detalle de enlaces ROTO (incluye soft-404)
    """
    col_m1, col_m2 = st.columns(2)
    col_m1.metric("✅ ACTIVOS", int((df_out["Status"] == "ACTIVO").sum()))
    col_m2.metric("❌ ROTOS", int((df_out["Status"] == "ROTO").sum()))

    df_view = df_out.copy()

    # Asegurar columna Nombre del Archivo
    if "Nombre del Archivo" not in df_view.columns and "Archivo" in df_view.columns:
        df_view = df_view.rename(columns={"Archivo": "Nombre del Archivo"})

    for col in ["Nombre del Archivo", "Link", "Status"]:
        if col not in df_view.columns:
            df_view[col] = ""

    df_view = df_view[["Nombre del Archivo", "Link", "Status"]]

    st.dataframe(
        _style_status_dataframe(df_view),
        use_container_width=True,
        height=420,
    )

    with st.expander("📊 Enlaces ROTOS (incluye soft-404)", expanded=False):
        rotos = df_out[df_out["Status"] == "ROTO"]
        if len(rotos) > 0:
            for _, row in rotos.iterrows():
                emoji = "🔴" if str(row.get("Soft_404", "")).strip().lower() == "sí" else "❌"
                score_info = f" (Score: {row.get('Score', 0)})"
                tipo = row.get("Tipo_Problema", "")
                st.markdown(
                    f"{emoji} Fila {row['Fila_Excel']}: `{row['Link']}`{score_info} → "
                    f"{row['Detalle']} ({tipo})"
                )
        else:
            st.success("No se detectaron enlaces con Status = ROTO.")

# ======= gestión robusta de módulo seleccionado =======================

def init_session_state():
    if "module" not in st.session_state:
        st.session_state["module"] = "Home"

    if "module_radio" not in st.session_state:
        st.session_state["module_radio"] = st.session_state["module"]

    st.session_state.setdefault("output_dir", str(Path.cwd() / "SALIDA_LINK_CHECKER"))
    st.session_state.setdefault("status_input_filename", None)
    st.session_state.setdefault("status_input_df", None)
    st.session_state.setdefault("status_links_list", None)
    st.session_state.setdefault("status_cache", {})
    st.session_state.setdefault("status_result_df", None)
    st.session_state.setdefault("status_invalid_df", None)
    st.session_state.setdefault("status_export_df", None)

    # Estado para PDF to Word Transform
    st.session_state.setdefault("extraccion_zip_bytes", None)
    st.session_state.setdefault("extraccion_resultados", None)
    st.session_state.setdefault("extraccion_errores", None)

    # Estado para Descarga Masiva
    st.session_state.setdefault("descarga_zip_bytes", None)
    st.session_state.setdefault("descarga_resultados", None)
    st.session_state.setdefault("descarga_fallidos", None)

    # Estado para Reporte Link (Word)
    st.session_state.setdefault("reporte_links_df", None)

    # 🔄 Estado del pipeline unificado "Report Broken Link"
    st.session_state.setdefault("pipeline_pdf_signature", None)
    st.session_state.setdefault("pipeline_pdf_done", False)
    st.session_state.setdefault("pipeline_pdf_results", None)
    st.session_state.setdefault("pipeline_pdf_errors", None)

    st.session_state.setdefault("pipeline_word_docs", None)
    st.session_state.setdefault("pipeline_ppt_docs", None)
    st.session_state.setdefault("pipeline_word_done", False)
    st.session_state.setdefault("pipeline_df_links", None)
    st.session_state.setdefault("pipeline_word_errors", None)

    st.session_state.setdefault("pipeline_word_inputs_count", 0) 
    st.session_state.setdefault("pipeline_ppt_inputs_count", 0)  

    st.session_state.setdefault("pipeline_status_done", False)

    # 🔁 Token para poder “reiniciar” el módulo unificado
    st.session_state.setdefault("pipeline_reset_token", 0)

def reset_report_broken_pipeline():
    """
    Limpia todo el estado relacionado al módulo unificado 'Report Broken Link'
    y fuerza que los widgets vuelvan a su estado inicial (como si se entrara
    por primera vez al módulo).
    """
    keys_to_clear = [
        # 🔹 1) Excel de Descarga Masiva (pasos 1–3)
        "pipeline_bulk_signature",
        "pipeline_bulk_done",
        "bulk_has_valid_urls",
        "bulk_urls_archivos",
        "descarga_resultados",
        "descarga_fallidos",
        "descarga_zip_bytes",

        # 🔹 2) PDF → Word (pasos 4–5)
        "pipeline_pdf_signature",
        "pipeline_pdf_done",
        "pipeline_pdf_results",
        "pipeline_pdf_errors",
        "extraccion_resultados",
        "extraccion_errores",
        "extraccion_zip_bytes",

        # Opciones de procesamiento de PDFs
        "extr_usar_multihilo",
        "extr_max_workers",

        # 🔹 3) Word → Links (pasos 6–7)
        "pipeline_word_docs",
        "pipeline_word_done",
        "pipeline_df_links",
        "pipeline_word_errors",
        "reporte_links_df",
        "pipeline_ppt_docs",
        "pipeline_word_inputs_count",
        "pipeline_ppt_inputs_count",


        # 🔹 4) Report Broken Link (validación de links, pasos 8–9)
        "pipeline_status_done",
        "status_input_filename",
        "status_input_df",
        "status_links_list",
        "status_cache",
        "status_result_df",
        "status_invalid_df",
        "status_export_df",
    ]

    for k in keys_to_clear:
        st.session_state.pop(k, None)

    # 🔁 Token para que file_uploader de Excel y PDF/ZIP se reinicien
    st.session_state["pipeline_reset_token"] = st.session_state.get(
        "pipeline_reset_token", 0
    ) + 1

def on_change_module():
    st.session_state["module"] = st.session_state["module_radio"]


# ======================================================
# DATA / EXPORT
# ======================================================
def _to_excel_report(df_status: pd.DataFrame) -> bytes:
    """
    Genera el Excel final SOLO con la hoja 'Status', con las columnas:

      name | Archivo | Página/Diapositiva | Link | Status | HTTP_Code | Detalle | Tipo_Problema | link_class

    Además, pinta la columna Status (ACTIVO=verde, ROTO=rojo).
    """
    from io import BytesIO
    try:
        from openpyxl.styles import PatternFill
    except ImportError:
        PatternFill = None  # Sin estilos si no está openpyxl.styles

    # ====== Preparar DataFrame para STATUS ======
    df_detalle = df_status.copy()

    # Asegurar nombre de columna Archivo
    if "Nombre del Archivo" in df_detalle.columns and "Archivo" not in df_detalle.columns:
        df_detalle = df_detalle.rename(columns={"Nombre del Archivo": "Archivo"})

    # Columnas objetivo en el orden requerido
    columnas_objetivo = [
        "name",
        "Archivo",
        "Página/Diapositiva",
        "Link",
        "Status",
        "HTTP_Code",
        "Detalle",
        "Tipo_Problema",
        "link_class",
    ]

    # Crear columnas faltantes si no existen
    for col in columnas_objetivo:
        if col not in df_detalle.columns:
            df_detalle[col] = ""

    # Reordenar columnas exactamente como se requiere
    df_detalle = df_detalle[columnas_objetivo]

    bio = BytesIO()
    with pd.ExcelWriter(bio, engine="openpyxl") as writer:
        # Hoja única: STATUS
        df_detalle.to_excel(writer, index=False, sheet_name="Status")

        if PatternFill is not None:
            # Colorear columna Status en la hoja Status
            green_fill = PatternFill(
                start_color="C6EFCE",
                end_color="C6EFCE",
                fill_type="solid",
            )
            red_fill = PatternFill(
                start_color="F8CBAD",
                end_color="F8CBAD",
                fill_type="solid",
            )

            ws = writer.sheets.get("Status")
            if ws is not None:
                # Buscar columna Status
                status_col_idx = None
                for cell in ws[1]:
                    if str(cell.value).strip().lower() == "status":
                        status_col_idx = cell.column
                        break

                if status_col_idx is not None:
                    for row_idx in range(2, ws.max_row + 1):
                        cell = ws.cell(row=row_idx, column=status_col_idx)
                        value = str(cell.value).strip().upper() if cell.value is not None else ""
                        if value == "ACTIVO":
                            cell.fill = green_fill
                        elif value == "ROTO":
                            cell.fill = red_fill

    bio.seek(0)
    return bio.getvalue()


def _to_excel_reporte_links(df_links: pd.DataFrame) -> bytes:
    """
    Genera un Excel con hoja 'Resultados' a partir del DataFrame de links.
    """
    from io import BytesIO

    bio = BytesIO()
    with pd.ExcelWriter(bio, engine="openpyxl") as writer:
        df_links.to_excel(writer, index=False, sheet_name="Resultados")
    return bio.getvalue()


def _read_excel_safe(uploaded_file) -> pd.DataFrame:
    try:
        return pd.read_excel(uploaded_file)
    except Exception as e:
        raise RuntimeError(f"No se pudo leer el Excel: {e}") from e


# ======================================================
# NORMALIZACIÓN DE URLs
# ======================================================

def _strip_invisible(s: str) -> str:
    return s.replace("\u200b", "").replace("\ufeff", "").strip()


def _looks_like_url(s: str) -> bool:
    s = s.lower().strip()
    return s.startswith(("http://", "https://")) or "." in s


def _normalize_one_url(
    raw: str,
    *,
    default_scheme: str = "https",
    allow_mailto: bool = False,
    allow_tel: bool = False,
    allow_anchors_only: bool = False,
) -> Tuple[Optional[str], str]:
    if raw is None:
        return None, "Vacío"

    s = _strip_invisible(str(raw))
    if not s:
        return None, "Vacío"

    s = s.replace("\n", " ").replace("\r", " ").strip()

    if s.startswith("#"):
        return (s, "") if allow_anchors_only else (None, "Anchor (#)")

    low = s.lower()
    if low.startswith("mailto:"):
        return (s, "") if allow_mailto else (None, "mailto")
    if low.startswith("tel:"):
        return (s, "") if allow_tel else (None, "tel")

    if not low.startswith(("http://", "https://")):
        if _looks_like_url(s):
            s = f"{default_scheme}://{s}"
        else:
            return None, "No parece URL"

    struct_valid, struct_reason = validate_url_structure(s)
    if not struct_valid:
        return None, struct_reason

    try:
        p = urlparse(s)
    except Exception:
        return None, "Parseo inválido"

    if not p.netloc:
        return None, "Sin dominio"

    netloc_raw = p.netloc.strip()
    userinfo = ""
    hostport = netloc_raw
    if "@" in netloc_raw:
        userinfo, hostport = netloc_raw.rsplit("@", 1)

    host = hostport
    port: Optional[str] = None
    if ":" in hostport:
        host, port = hostport.rsplit(":", 1)

    try:
        host_idn = host.encode("idna").decode("ascii")
    except Exception:
        host_idn = host

    scheme = p.scheme.lower()
    if port and (
        (scheme == "http" and port == "80")
        or (scheme == "https" and port == "443")
    ):
        port = None

    if port:
        netloc_clean = f"{host_idn}:{port}"
    else:
        netloc_clean = host_idn

    if userinfo:
        netloc_clean = f"{userinfo}@{netloc_clean}"

    path = quote(p.path, safe="/%:@-._~!$&'()*+,;=")
    query = quote(p.query, safe="=&%:@-._~!$&'()*+,;/?")

    norm = urlunparse((scheme, netloc_clean, path, p.params, query, ""))

    yt_valid, yt_reason = validate_youtube_url(norm)
    if not yt_valid:
        return None, yt_reason

    x_valid, x_reason = validate_x_twitter_url(norm)
    if not x_valid:
        return None, x_reason

    google_valid, google_reason = validate_google_search_url(norm)
    if not google_valid:
        return None, google_reason

    return norm, ""

def _normalize_links(
    series: pd.Series,
    *,
    allow_mailto: bool,
    allow_tel: bool,
    allow_anchors_only: bool,
    default_scheme: str,
) -> Tuple[List[Tuple[int, str]], pd.DataFrame]:
    out: List[Tuple[int, str]] = []
    invalid_rows: List[Dict[str, Any]] = []

    for excel_row, v in enumerate(series.tolist(), start=2):
        url, reason = _normalize_one_url(
            v,
            default_scheme=default_scheme,
            allow_mailto=allow_mailto,
            allow_tel=allow_tel,
            allow_anchors_only=allow_anchors_only,
        )
        if url is None:
            invalid_rows.append(
                {
                    "Fila_Excel": excel_row,
                    "Valor": "" if v is None else str(v),
                    "Motivo": reason,
                }
            )
            continue

        out.append((excel_row, url))

    return out, pd.DataFrame(invalid_rows)


# ======================================================
# LINK CHECKER ULTRA ROBUSTO V4
# ======================================================

def _httpx_available_or_warn() -> bool:
    if httpx is None:
        st.error("Falta la librería `httpx`. Instala con: `pip install httpx`")
        return False
    return True


def _requests_available_or_warn() -> bool:
    if requests is None:
        st.error("Falta la librería `requests`. Instala con: `pip install requests`")
        return False
    return True


# ======================================================
# CORE LINK CHECKER V5 (ULTRA ROBUSTO)
# ======================================================

def _host_key(url: str) -> str:
    try:
        p = urlparse(url)
        return p.netloc.lower()
    except Exception:
        return "unknown"


def _is_html_like(content_type: Optional[str]) -> bool:
    if not content_type:
        return False
    ct = content_type.lower()
    return "text/html" in ct or "application/xhtml" in ct

def _is_suspicious_redirect_to_root(original_url: str, final_url: str) -> bool:
    """
    Detecta el patrón típico de Google Sites en el que una URL profunda
    de contenido redirige a la raíz del sitio (home). Para el usuario,
    esto suele significar que la página concreta ya no existe.
    """
    try:
        p0 = urlparse(original_url)
        pf = urlparse(final_url)
    except Exception:
        return False

    domain = (pf.netloc or "").lower()

    # De momento, aplicamos la heurística sólo a Google Sites
    if "sites.google.com" not in domain:
        return False

    seg0 = [s for s in (p0.path or "").split("/") if s]
    segf = [s for s in (pf.path or "").split("/") if s]

    # La URL original debe ser claramente "más profunda"
    if len(seg0) <= len(segf):
        return False

    # La ruta final debe ser prefijo de la inicial
    if seg0[: len(segf)] != segf:
        return False

    # Consideramos sospechoso cuando se redirige a la raíz o a 'home'
    #   /site/<sitio>/
    #   /site/<sitio>/home
    if len(segf) <= 2:
        return True
    if len(segf) == 3 and segf[-1].lower() in {"home", "inicio", "index", "default"}:
        return True

    return False


def _soft_404_detect_v5(body_text: str, url: str) -> Tuple[bool, int]:
    """
    Detección reforzada de soft-404:

    1) Primero aplica patrones fuertes (SOFT_404_STRONG_RE) que indican
       claramente que el recurso no existe o no es accesible:
       - YouTube: "Video no disponible", "Este video es privado", etc.
       - X: "Esta página no existe. Intenta hacer otra búsqueda."
       - Otros mensajes típicos de "page not found".

    2) Si no hay patrón fuerte, usa el score heurístico de
       `_calculate_content_score`. Un score muy negativo se interpreta
       como soft-404 genérico.
    """
    if not body_text:
        return False, 0

    # Trozo suficiente para buscar mensajes de error
    chunk = body_text[:8000]

    # 1) Patrones fuertes → confianza muy alta, en cualquier dominio
    if SOFT_404_STRONG_RE.search(chunk):
        return True, 95

    # 2) Score heurístico (penaliza SOFT_404_RE, textos muy cortos, etc.)
    score = _calculate_content_score(body_text, url)

    # Score muy negativo => consideramos soft-404
    if score < -10:
        return True, min(90, abs(score))

    return False, 0

def _classify_v5(
    url: str,
    status_code: Optional[int],
    detail: str,
    redirected: bool,
) -> str:
    """
    Clasificación mejorada:
    - Usa _is_valid_status_code para considerar ciertos 4xx (Facebook, Canva) como ACTIVO.
    - Mantiene la lógica original para 404/410 y 5xx (ERROR_SERVIDOR vía Tipo_Problema).
    """
    if status_code is None:
        return "ERROR"

    # Si para este dominio el código se considera válido, lo tratamos como ACTIVO/REDIRECT
    if _is_valid_status_code(url, status_code):
        return "REDIRECT" if redirected else "ACTIVO"

    # No válido → aplicamos regla original
    if status_code in (404, 410):
        return "ROTO"

    if 500 <= status_code <= 599:
        return "ERROR"

    if 400 <= status_code <= 499:
        return "ERROR"

    return "ERROR"


def _is_binary_candidate(url: str) -> bool:
    try:
        p = urlparse(url)
        path = p.path.lower()
    except Exception:
        return False
    return any(path.endswith(ext) for ext in BINARY_EXTS)


def _compute_retry_delay(retry_after_header: Optional[str], attempt: int) -> float:
    if retry_after_header:
        try:
            return float(retry_after_header)
        except Exception:
            pass
    # backoff exponencial básico + algo de jitter
    return min(30, 1.0 * (2 ** (attempt - 1))) + random.random()


async def _fetch_limited_text_v5(
    client: "httpx.AsyncClient",
    url: str,
    timeout_s: float,
    max_bytes: int,
    range_bytes: int,
) -> Tuple[Optional[int], Dict[str, str], str, bool, str, List[str]]:
    """
    GET parcial con cabecera Range. Se limita el contenido a max_bytes.
    Combina los headers del cliente con Range.
    """
    headers = {"Range": f"bytes=0-{range_bytes-1}"}

    try:
        async with client.stream(
            "GET",
            url,
            timeout=timeout_s,
            follow_redirects=True,
            headers=headers,
        ) as r:
            final_url = str(r.url)
            history_urls = [str(resp.url) for resp in r.history]
            redirect_chain = (
                history_urls + [final_url]
                if history_urls or final_url != url
                else [final_url]
            )
            redirected = final_url != url or bool(history_urls)
            status = r.status_code
            h = {k.lower(): v for k, v in r.headers.items()}

            buf = bytearray()
            async for chunk in r.aiter_bytes():
                if not chunk:
                    continue
                take = min(len(chunk), max_bytes - len(buf))
                buf.extend(chunk[:take])
                if len(buf) >= max_bytes:
                    break

            encoding = r.encoding or "utf-8"
            try:
                text = buf.decode(encoding, errors="replace")
            except Exception:
                text = buf.decode("utf-8", errors="replace")

            return status, h, text, redirected, final_url, redirect_chain

    except Exception as e:
        return None, {}, f"{e.__class__.__name__}: {str(e)[:200]}", False, url, [url]


async def _check_one_url_robust_v5(
    client: "httpx.AsyncClient",
    url: str,
    *,
    timeout_s: float,
    max_bytes: int,
    range_bytes: int,
    detect_soft_404: bool,
    retries: int,
) -> Dict[str, Any]:
    """
    Verificación robusta V5:
    - HEAD para binarios (PDF, DOC, etc.).
    - GET parcial para HTML.
    - User-Agent realista + headers por dominio (aportados por el cliente).
    - Soft-404 con scoring y dominios confiables.
    - Reglas específicas por dominio (SSL, códigos aceptados, Google Sites, etc.).
    - 🔴 Regla especial: enlaces que contienen 'canvas.utp' se marcan SIEMPRE como ROTO.
    """
    now_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    # 🔴 0) Regla especial Canvas UTP: cualquier enlace que contenga 'canvas.utp'
    # se marca directamente como ROTO, sin hacer ninguna petición HTTP.
    # Esto aplica tanto si el enlace está accesible como si no: para el reporte
    # institucional se consideran rotos por definición.
    if CANVAS_UTP_KEYWORD in url.lower():
        return {
            "Link": url,
            "Status": "ROTO",
            "HTTP_Code": None,  # No se consulta HTTP
            "Detalle": "Dominio canvas.utp marcado como ROTO por regla institucional UTP",
            "Content_Type": "",
            "Redirected": "No",
            "Timestamp": now_str,
            "Final_URL": url,
            "Redirect_Chain": url,
            "Soft_404": "No",
            "Score": -100,
        }

    # 1) HEAD para binarios (PDF, DOC, etc.)
    if _is_binary_candidate(url):
        attempt = 0
        while attempt <= max(0, retries):
            attempt += 1
            try:
                r = await client.head(
                    url,
                    timeout=timeout_s,
                    follow_redirects=True,
                )
                final_url = str(r.url)
                history_urls = [str(resp.url) for resp in r.history]
                redirect_chain = (
                    history_urls + [final_url]
                    if history_urls or final_url != url
                    else [final_url]
                )
                redirected = final_url != url or bool(history_urls)

                status = r.status_code
                headers = {k.lower(): v for k, v in r.headers.items()}
                ct = headers.get("content-type", "")

                # Algunos servidores no soportan HEAD correctamente
                if status in (405, 501):
                    break

                # Códigos transitorios → reintentos
                if status in (408, 425, 429) or (500 <= status <= 599):
                    delay = _compute_retry_delay(headers.get("retry-after"), attempt)
                    if attempt <= retries:
                        await asyncio.sleep(delay)
                        continue

                ct_valid, ct_reason = validate_content_type_match(url, ct)
                if not ct_valid:
                    return {
                        "Link": url,
                        "Status": "ROTO",
                        "HTTP_Code": status,
                        "Detalle": f"Content-Type incorrecto: {ct_reason}",
                        "Content_Type": ct,
                        "Redirected": "Sí" if redirected else "No",
                        "Timestamp": now_str,
                        "Final_URL": final_url,
                        "Redirect_Chain": " -> ".join(redirect_chain),
                        "Soft_404": "Sí",
                        "Score": -50,
                    }

                detail = "OK" if _is_valid_status_code(url, status) else f"HTTP {status}"

                return {
                    "Link": url,
                    "Status": _classify_v5(url, status, detail, redirected),
                    "HTTP_Code": status,
                    "Detalle": detail,
                    "Content_Type": ct,
                    "Redirected": "Sí" if redirected else "No",
                    "Timestamp": now_str,
                    "Final_URL": final_url,
                    "Redirect_Chain": " -> ".join(redirect_chain),
                    "Soft_404": "No",
                    "Score": 100,
                }

            except Exception as e:
                last_detail = f"{e.__class__.__name__}: {str(e)[:200]}"
                if attempt <= retries:
                    delay = _compute_retry_delay(None, attempt)
                    await asyncio.sleep(delay)
                    continue

                return {
                    "Link": url,
                    "Status": "ERROR",
                    "HTTP_Code": None,
                    "Detalle": last_detail,
                    "Content_Type": "",
                    "Redirected": "No",
                    "Timestamp": now_str,
                    "Final_URL": url,
                    "Redirect_Chain": url,
                    "Soft_404": "No",
                    "Score": -100,
                }

    # 2) GET parcial para HTML / otros
    attempt = 0
    last_detail = ""
    last_status: Optional[int] = None
    last_redirected = False
    last_ct = ""
    last_final_url = url
    last_chain: List[str] = [url]
    soft_flag = False
    content_score = 0

    while attempt <= max(0, retries):
        attempt += 1

        status, headers, text, redirected, final_url, chain = await _fetch_limited_text_v5(
            client,
            url,
            timeout_s=timeout_s,
            max_bytes=max_bytes,
            range_bytes=range_bytes,
        )

        last_status = status
        last_redirected = redirected
        last_final_url = final_url
        last_chain = chain
        last_ct = headers.get("content-type", "")

        # Error de red / timeout sin status
        if status is None:
            last_detail = text
            if attempt <= retries:
                delay = _compute_retry_delay(None, attempt)
                await asyncio.sleep(delay)
                continue
            break

        # Códigos transitorios → reintentos
        if status in (408, 425, 429) or (500 <= status <= 599):
            last_detail = f"HTTP {status} (transitorio)"
            if attempt <= retries:
                delay = _compute_retry_delay(headers.get("retry-after"), attempt)
                await asyncio.sleep(delay)
                continue
            break

        # Códigos >= 400 no aceptados por dominio → error
        if not _is_valid_status_code(url, status):
            if status >= 400:
                last_detail = f"HTTP {status}"
                break

        ct_valid, ct_reason = validate_content_type_match(url, last_ct)
        if not ct_valid:
            soft_flag = True
            last_detail = f"Content-Type incorrecto: {ct_reason}"
            return {
                "Link": url,
                "Status": "ROTO",
                "HTTP_Code": status,
                "Detalle": last_detail,
                "Content_Type": last_ct,
                "Redirected": "Sí" if redirected else "No",
                "Timestamp": now_str,
                "Final_URL": final_url,
                "Redirect_Chain": " -> ".join(chain),
                "Soft_404": "Sí",
                "Score": -50,
            }

        detail = "OK"

        # Soft-404 con scoring (sólo si es HTML)
        if detect_soft_404 and _is_html_like(last_ct):
            is_soft, conf = _soft_404_detect_v5(text, final_url)
            content_score = _calculate_content_score(text, final_url)

            if is_soft:
                soft_flag = True
                last_detail = (
                    f"Soft-404 detectado (confianza: {conf}%, score: {content_score})"
                )
                return {
                    "Link": url,
                    "Status": "ROTO",
                    "HTTP_Code": status,
                    "Detalle": last_detail,
                    "Content_Type": last_ct,
                    "Redirected": "Sí" if redirected else "No",
                    "Timestamp": now_str,
                    "Final_URL": final_url,
                    "Redirect_Chain": " -> ".join(chain),
                    "Soft_404": "Sí",
                    "Score": content_score,
                }

        # 🔴 Heurística para redirección sospechosa a la raíz (Google Sites)
        if redirected and status is not None and status < 400:
            if _is_suspicious_redirect_to_root(url, final_url):
                soft_flag = True
                last_detail = (
                    "Redirección a la raíz del sitio (home). "
                    "Probablemente la página específica ya no existe (soft-404 por redirect)."
                )
                return {
                    "Link": url,
                    "Status": "ROTO",
                    "HTTP_Code": status,
                    "Detalle": last_detail,
                    "Content_Type": last_ct,
                    "Redirected": "Sí",
                    "Timestamp": now_str,
                    "Final_URL": final_url,
                    "Redirect_Chain": " -> ".join(chain),
                    "Soft_404": "Sí",
                    "Score": -25,
                }

        # Si llegamos aquí, lo consideramos válido (ACTIVO / REDIRECT)
        return {
            "Link": url,
            "Status": _classify_v5(url, status, detail, redirected),
            "HTTP_Code": status,
            "Detalle": detail,
            "Content_Type": last_ct,
            "Redirected": "Sí" if redirected else "No",
            "Timestamp": now_str,
            "Final_URL": final_url,
            "Redirect_Chain": " -> ".join(chain),
            "Soft_404": "No",
            "Score": content_score if content_score else 100,
        }

    # Caso de salida por error tras reintentos
    return {
        "Link": url,
        "Status": _classify_v5(url, last_status, last_detail or "Error", last_redirected),
        "HTTP_Code": last_status,
        "Detalle": last_detail or "Error",
        "Content_Type": last_ct,
        "Redirected": "Sí" if last_redirected else "No",
        "Timestamp": now_str,
        "Final_URL": last_final_url,
        "Redirect_Chain": " -> ".join(last_chain),
        "Soft_404": "Sí" if soft_flag else "No",
        "Score": content_score,
    }

async def _run_link_check_ultra_v5(
    links_with_rows: List[Tuple[int, str]],
    *,
    timeout_s: float,
    concurrency_global: int,
    concurrency_per_host: int,
    detect_soft_404: bool,
    retries: int,
    verify_ssl: bool,
    max_bytes: int,
    range_bytes: int,
    progress_callback,
) -> List[Dict[str, Any]]:
    """
    Versión V5 del checker:
    - Concurrencia global y por host.
    - Cache por URL en session_state["status_cache"].
    - Headers por dominio (User-Agent realista).
    - **NUEVO**: dos clientes HTTPX
        * client_ssl   → verify = verify_ssl (según toggle)
        * client_nossl → verify = False (para dominios con skip_ssl_verify)
    """
    sem_global = asyncio.Semaphore(max(1, int(concurrency_global)))
    host_sems: Dict[str, asyncio.Semaphore] = {}

    def get_host_sem(url: str) -> asyncio.Semaphore:
        hk = _host_key(url)
        if hk not in host_sems:
            host_sems[hk] = asyncio.Semaphore(max(1, int(concurrency_per_host)))
        return host_sems[hk]

    limits = httpx.Limits(
        max_connections=max(10, int(concurrency_global) + 10),
        max_keepalive_connections=max(10, int(concurrency_global)),
        keepalive_expiry=30.0,
    )

    timeout = httpx.Timeout(timeout_s)

    cache: Dict[str, Dict[str, Any]] = st.session_state.get("status_cache", {})
    if not isinstance(cache, dict):
        cache = {}
    st.session_state["status_cache"] = cache

    # Headers base genéricos
    base_headers = _build_headers_for_domain("https://example.com")

    # 👇 Creamos dos clientes: uno respeta el toggle verify_ssl, el otro siempre ignora certificados
    async with httpx.AsyncClient(
        headers=base_headers,
        limits=limits,
        timeout=timeout,
        http2=False,
        verify=verify_ssl,
        follow_redirects=True,
    ) as client_ssl, httpx.AsyncClient(
        headers=base_headers,
        limits=limits,
        timeout=timeout,
        http2=False,
        verify=False,
        follow_redirects=True,
    ) as client_nossl:

        total = len(links_with_rows)
        done = 0
        results: List[Dict[str, Any]] = []

        async def worker(fila_excel: int, u: str):
            nonlocal done
            host_sem = get_host_sem(u)

            # Config por dominio
            config = _get_domain_config(u)

            # Si el usuario desactivó "Verificar SSL" → usamos siempre client_nossl.
            # Si está activado, sólo usamos client_nossl para dominios que tengan skip_ssl_verify = True
            use_no_ssl = (not verify_ssl) or bool(config.get("skip_ssl_verify", False))
            client = client_nossl if use_no_ssl else client_ssl

            # Headers ajustados por dominio (User-Agent, Referer, etc.)
            client.headers.update(_build_headers_for_domain(u))

            async with sem_global:
                async with host_sem:
                    if u in cache:
                        base = cache[u]
                    else:
                        base = await _check_one_url_robust_v5(
                            client,
                            u,
                            timeout_s=timeout_s,
                            max_bytes=max_bytes,
                            range_bytes=range_bytes,
                            detect_soft_404=detect_soft_404,
                            retries=retries,
                        )
                        cache[u] = base

            row = dict(base)
            row["Fila_Excel"] = fila_excel

            done += 1
            progress_callback(done, total, u, row.get("Status", ""))
            return row

        tasks = [worker(fila, url) for (fila, url) in links_with_rows]
        for coro in asyncio.as_completed(tasks):
            results.append(await coro)

        st.session_state["status_cache"] = cache
        return results

def run_async(coro):
    """
    Wrapper compatible con Streamlit (maneja el caso de event loop ya existente).
    """
    try:
        return asyncio.run(coro)
    except RuntimeError as e:
        msg = str(e)
        if "asyncio.run() cannot be called from a running event loop" in msg:
            loop = asyncio.get_event_loop()
            return loop.run_until_complete(coro)
        raise

# ======================================================
# HELPER PARA TIPO_PROBLEMA
# ======================================================

def _infer_tipo_problema(row: pd.Series) -> str:
    """
    Devuelve una etiqueta más específica del problema:
    - SIN_PROBLEMA
    - SOFT_404
    - ROTO_REAL
    - ACCESO_RESTRINGIDO
    - ERROR_SERVIDOR
    - ERROR_CLIENTE
    - ERROR_DESCONOCIDO
    - FORMATO_INVALIDO
    """
    status = str(row.get("Status", "") or "").upper()
    code = row.get("HTTP_Code")
    detalle = str(row.get("Detalle", "") or "")
    soft_404_flag = str(row.get("Soft_404", "") or "").strip().lower() == "sí"

    # Normalizar code a int o None
    try:
        if pd.isna(code):
            code = None
    except Exception:
        pass
    if isinstance(code, float):
        try:
            code = int(code)
        except Exception:
            pass

    # INVALIDO por estructura de URL
    if status == "INVALIDO":
        return "FORMATO_INVALIDO"

    # Activo / redirect -> sin problema
    if status in ("ACTIVO", "REDIRECT"):
        return "SIN_PROBLEMA"

    # ROTO (404/410 o soft-404 explícito / Content-Type incorrecto)
    if status == "ROTO":
        if soft_404_flag or "soft-404" in detalle.lower() or "content-type incorrecto" in detalle.lower():
            return "SOFT_404"
        if code in (404, 410):
            return "ROTO_REAL"
        return "ROTO_REAL"

    # ERROR (resto 4xx/5xx o errores de red)
    if status == "ERROR":
        if code in (401, 403, 429):
            return "ACCESO_RESTRINGIDO"
        if code is None:
            return "ERROR_DESCONOCIDO"
        try:
            c_int = int(code)
        except Exception:
            return "ERROR_DESCONOCIDO"

        if 500 <= c_int <= 599:
            return "ERROR_SERVIDOR"
        if 400 <= c_int <= 499:
            return "ERROR_CLIENTE"
        return "ERROR_DESCONOCIDO"

    return ""

def _standardize_status_column(df: pd.DataFrame) -> pd.DataFrame:
    """
    Normaliza la columna Status según la regla:
    - REDIRECT → ACTIVO
    - ERROR / INVALIDO → ROTO
    Mantiene otros valores (como ACTIVO, ROTO) sin cambios.
    """
    if "Status" not in df.columns:
        return df

    df = df.copy()
    status_upper = df["Status"].astype(str).str.upper()

    # REDIRECT se considera ACTIVO
    df.loc[status_upper.str.contains("REDIRECT"), "Status"] = "ACTIVO"

    # ERROR o INVALIDO se consideran ROTO
    df.loc[
        status_upper.str.contains("ERROR") | status_upper.str.contains("INVALIDO"),
        "Status"
    ] = "ROTO"

    return df

# Prefijo automático generado por el módulo Descarga Masiva:
# Descarga_Masiva_Documentos_YYYYMMDD_HHMMSS_<nombre_original>.pdf
DESCARGA_PREFIX_RE = re.compile(
    r"^Descarga_Masiva_Documentos_\d{8}_\d{6}_",
    re.IGNORECASE,
)

# ======================================================
# PDF to Word Transform - PDF → WORD (LÓGICA DEL EXTRACTOR)
# ======================================================

BIB_HEAD_PATTERNS = [
    "fuentes bibliograficas",
    "referencias bibliograficas",
    "referenciasbibliograficas",
    "bibliografia",
    "referencias",
    "obras citadas",
    "citas bibliograficas",
    "webgrafia",
]

RE_APA_YEAR = re.compile(r"\(\d{4}[a-z]?\)")
RE_YEAR_AFTER_DOT = re.compile(r"\b\d{4}\.")
RE_DOI = re.compile(r"\b10\.\d{4,9}/\S+\b", re.IGNORECASE)
RE_URL = re.compile(r"https?://\S+|www\.\S+", re.IGNORECASE)
RE_ISBN = re.compile(r"\bISBN(?:-13)?:?\s*(97(8|9))?[\d\- ]{8,}", re.IGNORECASE)
RE_ISSN = re.compile(r"\bISSN:?[\d\- ]{7,}", re.IGNORECASE)
RE_CORCHETES_NUM = re.compile(r"^\s*\[\d+([,\-]\d+)*\]\s*")
RE_BULLET = re.compile(r"^\s*(?:•|-|\d+[\.\)]|\([a-zA-Z0-9]\))\s+")
RE_FUENTE_HEADER = re.compile(r"^\s*fuente[s]?\s*:\s*$", re.IGNORECASE)
RE_FUENTE_INLINE = re.compile(r"^\s*fuente[s]?\s*:", re.IGNORECASE)

class PDFBatchProcessor:
    """Procesador optimizado para múltiples archivos PDF (lógica original, sin GUI Tk)."""

    def __init__(self, max_workers: int = None):
        self.max_workers = max_workers or min(4, os.cpu_count() or 1)
        self.cancel_event = threading.Event()
        self.progress_queue: "queue.Queue[Dict[str, Any]]" = queue.Queue()

    def process_single_pdf(self, pdf_path: str, output_dir: str, options: Dict) -> Dict:
        """Procesa un único archivo PDF → Word (permite filtrar o no bibliografía según `options`)."""
        if self.cancel_event.is_set():
            return {"status": "cancelled", "file": pdf_path}

        if fitz is None or Document is None:
            return {
                "status": "error",
                "file": pdf_path,
                "error": "Faltan dependencias `pymupdf` o `python-docx`.",
                "elapsed_time": 0,
                "success": False,
            }

        try:
            start_time = datetime.now()
            pdf_name = Path(pdf_path).name

            self.progress_queue.put({"type": "file_start", "file": pdf_name, "total_pages": 0})

            success, output_path, stats = self._process_pdf_internal(pdf_path, output_dir, options)

            elapsed = (datetime.now() - start_time).total_seconds()

            result: Dict[str, Any] = {
                "status": "success" if success else "error",
                "file": pdf_path,
                "output": output_path,
                "stats": stats,
                "elapsed_time": elapsed,
                "success": success,
            }

            if not success:
                result["error"] = stats.get("error", "Error desconocido")

            return result

        except Exception as e:
            logger.error(f"Error procesando {pdf_path}: {e}")
            return {
                "status": "error",
                "file": pdf_path,
                "error": str(e),
                "elapsed_time": 0,
                "success": False,
            }

    def _process_pdf_internal(self, pdf_path: str, output_dir: str, options: Dict) -> Tuple[bool, str, Dict]:
        """Lógica interna de procesamiento de PDF."""
        use_multithread = options.get("usar_multihilo", True)

        try:
            with fitz.open(pdf_path) as doc:
                num_pages = len(doc)

            self.progress_queue.put(
                {"type": "file_pages", "file": Path(pdf_path).name, "total_pages": num_pages}
            )

            output_path = Path(output_dir) / f"{Path(pdf_path).stem}.docx"
            doc_word = Document()

            if use_multithread and num_pages > 3:
                from concurrent.futures import ThreadPoolExecutor

                futures: List[Tuple[int, Any]] = []
                results: List[Tuple[int, str]] = []
                with ThreadPoolExecutor(max_workers=options.get("max_workers", 4)) as executor:
                    for page_num in range(num_pages):
                        if self.cancel_event.is_set():
                            break
                        future = executor.submit(
                            self._extract_and_process_page,
                            pdf_path,
                            page_num,
                            options,
                        )
                        futures.append((page_num, future))

                    for page_num, future in futures:
                        if self.cancel_event.is_set():
                            break
                        page_idx, text = future.result()
                        results.append((page_idx, text))

                        self.progress_queue.put(
                            {
                                "type": "page_progress",
                                "file": Path(pdf_path).name,
                                "page": page_num + 1,
                                "total": num_pages,
                            }
                        )

                results.sort(key=lambda x: x[0])
                page_texts = [text for _, text in results]
            else:
                page_texts: List[str] = []
                for page_num in range(num_pages):
                    if self.cancel_event.is_set():
                        break
                    _, text = self._extract_and_process_page(pdf_path, page_num, options)
                    page_texts.append(text)

                    self.progress_queue.put(
                        {
                            "type": "page_progress",
                            "file": Path(pdf_path).name,
                            "page": page_num + 1,
                            "total": num_pages,
                        }
                    )

            if self.cancel_event.is_set():
                return False, str(output_path), {"status": "cancelled"}

            # Volcar el texto procesado al documento Word
            for idx, text in enumerate(page_texts):
                if text.strip():
                    for line in text.split("\n"):
                        if line.strip():
                            doc_word.add_paragraph(line)
                if idx < len(page_texts) - 1:
                    doc_word.add_page_break()

            doc_word.save(str(output_path))

            stats = {
                "archivo": pdf_path,
                "nombre_archivo": Path(pdf_path).stem,
                "paginas_procesadas": num_pages,
                "archivo_salida": str(output_path),
                "tamano_salida": os.path.getsize(output_path) if os.path.exists(output_path) else 0,
                "errores": sum(1 for text in page_texts if "ERROR:" in text),
                "timestamp": datetime.now().isoformat(),
            }

            return True, str(output_path), stats

        except Exception as e:
            logger.error(f"Error interno procesando {pdf_path}: {e}")
            return False, str(Path(output_dir) / f"{Path(pdf_path).stem}.docx"), {"error": str(e)}

    def _extract_and_process_page(self, pdf_path: str, page_num: int, options: Dict) -> Tuple[int, str]:
        """Extrae y procesa una página individual."""
        try:
            with fitz.open(pdf_path) as doc:
                if page_num >= len(doc):
                    return page_num, f"=== PÁGINA {page_num + 1} ===\nERROR: Página no existe"

                page = doc[page_num]
                raw_text = page.get_text()

                if len(raw_text.strip()) < 50:
                    raw_text = page.get_text("text")

                cleaned_text = self._clean_text(raw_text)

                # Solo se filtran referencias si la opción lo indica explícitamente
                if options.get("filtrar_bibliografia", False):
                    text_base = self._filter_references(cleaned_text)
                else:
                    text_base = cleaned_text

                text_no_formulas = self._filter_formulas(text_base)
                reformatted_text = self._reformat_sentences(text_no_formulas)

                result = f"=== PÁGINA {page_num + 1} ===\n{reformatted_text.strip()}"
                return page_num, result

        except Exception as e:
            error_msg = f"ERROR procesando página {page_num + 1}: {str(e)}"
            logger.error(f"{error_msg} en {pdf_path}")
            return page_num, f"=== PÁGINA {page_num + 1} ===\n{error_msg}"

    def _clean_text(self, text: str) -> str:
        """Limpia el texto (normalización básica + unión de líneas)."""
        replacements = {
            "\x00": "",
            "\x0c": "\n",
            "\uf0b7": "•",
            "\uf0a7": "§",
            "\uf0d8": "°",
            "\xad": "",
            "\t": "    ",
        }

        for old, new in replacements.items():
            text = text.replace(old, new)

        lines = [line.strip() for line in text.split("\n")]
        lines = [line for line in lines if line]

        final_lines: List[str] = []
        buffer: List[str] = []

        def flush_buffer():
            nonlocal buffer, final_lines
            if buffer:
                final_lines.append(" ".join(buffer))
                buffer = []

        for line in lines:
            # Si parece ecuación o línea de referencia, no la unimos con el buffer
            if "=" in line or self._is_reference_line(line):
                flush_buffer()
                final_lines.append(line)
                continue

            if len(line) < 80 and not line.endswith((".", "!", "?", ":", ";", ",", ")")):
                buffer.append(line)
            else:
                if buffer:
                    final_lines.append(" ".join(buffer + [line]))
                    buffer = []
                else:
                    final_lines.append(line)

        flush_buffer()
        return "\n".join(final_lines).strip()

    def _is_reference_line(self, line: str) -> bool:
        """Determina si una línea es una referencia bibliográfica."""
        if not line or len(line.strip()) < 5:
            return False

        text = line.strip()

        if (
            RE_APA_YEAR.search(text)
            or RE_YEAR_AFTER_DOT.search(text)
            or RE_DOI.search(text)
            or RE_URL.search(text)
            or RE_ISBN.search(text)
            or RE_ISSN.search(text)
            or RE_CORCHETES_NUM.search(text)
        ):
            return True

        if RE_BULLET.match(text):
            if len(text.split()) >= 3:
                return True

        numbers = sum(ch.isdigit() for ch in text)
        if numbers >= 4 and ("," in text or "." in text):
            if "pp." in text.lower() or "p." in text.lower():
                return True

        return False

    def _filter_references(self, text: str) -> str:
        """Filtra bloques de referencias / bibliografía."""
        if not text.strip():
            return text

        lines = text.split("\n")
        result: List[str] = []
        in_ref_block = False

        for i, line in enumerate(lines):
            norm = self._normalize_text(line)

            if in_ref_block:
                if not line.strip():
                    continue
                if self._is_reference_line(line):
                    continue
                in_ref_block = False

            is_bib_header = any(pattern in norm for pattern in BIB_HEAD_PATTERNS)

            if is_bib_header and self._is_reference_header(lines, i):
                in_ref_block = True
                continue

            if RE_FUENTE_INLINE.match(line):
                if RE_FUENTE_HEADER.match(line):
                    ref_count = 0
                    for j in range(i + 1, min(len(lines), i + 6)):
                        if self._is_reference_line(lines[j]):
                            ref_count += 1
                    if ref_count > 0:
                        in_ref_block = True
                        continue
                else:
                    # “Fuente: …” en la misma línea → se asume referencia y se omite
                    continue

            result.append(line)

        return "\n".join(result).strip()

    def _filter_formulas(self, text: str) -> str:
        """En esta versión no se eliminan fórmulas; se devuelve el texto tal cual."""
        return text

    def _reformat_sentences(self, text: str) -> str:
        """Reformatea el texto para dejar una oración por línea."""
        if not text.strip():
            return text

        # Unir todo en una sola línea base
        text = re.sub(r"\s*\n\s*", " ", text)
        text = re.sub(r"\s+", " ", text).strip()

        # Proteger números tipo 3.14 para que no se separen mal
        text = re.sub(r"(\d+)\.\s+(\d{1,3})", r"\1.\2", text)

        def protect_dots(match):
            return match.group(0).replace(".", "[[DOT_PAREN]]")

        text = re.sub(r"\([^)]*\)", protect_dots, text)

        # Proteger numeraciones tipo "1. Introducción"
        text = re.sub(r"\b(\d+)\.\s+(?=[A-ZÁÉÍÓÚÑ])", r"\1§ ", text)

        # Cortar oraciones en puntos seguidos
        text = re.sub(r"(?<!\d)\.\s+(?!\d)", ".\n", text)
        text = re.sub(r"\)\s+(?=[A-ZÁÉÍÓÚÑ¿])", ")\n", text)

        # Restaurar marcadores
        text = text.replace("§", ".")
        text = text.replace("[[DOT_PAREN]]", ".")

        lines = [line.strip() for line in text.split("\n")]
        lines = [line for line in lines if line]

        return "\n".join(lines)

    def _normalize_text(self, text: str) -> str:
        """Normaliza texto para comparaciones (sin tildes, minúsculas, espacios colapsados)."""
        if not text:
            return ""
        text = text.lower()
        text = "".join(
            c for c in unicodedata.normalize("NFD", text)
            if unicodedata.category(c) != "Mn"
        )
        text = re.sub(r"\s+", " ", text).strip()
        return text

    def _is_reference_header(self, lines: List[str], idx: int) -> bool:
        """Determina si una línea es cabecera de la sección de referencias."""
        if idx < 0 or idx >= len(lines):
            return False

        start = idx + 1
        end = min(len(lines), idx + 11)

        window = [line for line in lines[start:end] if line.strip()]
        if not window:
            return False

        total = len(window)
        ref_like = sum(1 for line in window if self._is_reference_line(line))

        return ref_like >= 2 or ref_like / total >= 0.4

# ======================================================
# DESCARGA MASIVA (DOCS PDF/PPT/WORD)
# ======================================================

def _format_hms(seconds: float) -> str:
    """Formatea segundos a HH:MM:SS o MM:SS."""
    seconds_int = int(max(0, seconds))
    m, s = divmod(seconds_int, 60)
    h, m = divmod(m, 60)
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"


def nombre_archivo_seguro(url: str, carpeta_destino: str, max_ruta: int = 240) -> str:
    """
    Genera un nombre de archivo seguro para Windows:
    - Decodifica caracteres %xx
    - Elimina querystrings
    - Reemplaza caracteres inválidos
    - Acorta el nombre si la ruta se hace demasiado larga
    """
    nombre = url.split('/')[-1]
    nombre = nombre.split('?')[0]
    from urllib.parse import unquote as _unq
    nombre = _unq(nombre)

    caracteres_invalidos = '<>:"/\\|?*'
    for c in caracteres_invalidos:
        nombre = nombre.replace(c, '_')

    if not nombre.strip():
        nombre = "archivo_descargado"

    ruta_base = os.path.join(carpeta_destino, "")
    espacio_disponible = max_ruta - len(ruta_base)
    if espacio_disponible < 50:
        espacio_disponible = 50

    if len(nombre) > espacio_disponible:
        base, ext = os.path.splitext(nombre)
        if len(ext) > 10:
            ext = ext[:10]
        max_base = espacio_disponible - len(ext)
        if max_base < 1:
            max_base = 1
        base = base[:max_base]
        nombre = base + ext

    return nombre


def _run_descarga_masiva_streamlit(
    urls_archivos: List[str],
    *,
    progress_bar,
    progress_text,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], Optional[bytes]]:
    """
    Ejecuta la descarga masiva de documentos en un directorio temporal,
    mostrando progreso en Streamlit y devolviendo:
      - resultados OK
      - fallidos
      - zip_bytes con todos los archivos descargados + CSV de fallidos (si aplica)
    """
    if not urls_archivos:
        return [], [], None

    tmp_dir = Path(tempfile.mkdtemp(prefix="utp_descarga_masiva_"))
    session = requests.Session()

    resultados: List[Dict[str, Any]] = []
    fallidos: List[Dict[str, Any]] = []
    total = len(urls_archivos)
    start_time = time.time()

    for idx, url in enumerate(urls_archivos, start=1):
        url = str(url).strip()
        descargado_ok = False
        ultimo_error = ""
        nombre_archivo = ""
        ruta_archivo = None

        try:
            nombre_archivo = nombre_archivo_seguro(url, str(tmp_dir))
            ruta_archivo = tmp_dir / nombre_archivo

            if ruta_archivo.exists():
                base, ext = os.path.splitext(nombre_archivo)
                contador = 1
                while ruta_archivo.exists():
                    sufijo = f"_{contador}"
                    espacio_disponible = 240 - len(os.path.join(str(tmp_dir), ""))
                    max_base = espacio_disponible - len(ext) - len(sufijo)
                    if max_base < 1:
                        max_base = 1
                    base_trunc = base[:max_base]
                    nombre_archivo_alt = f"{base_trunc}{sufijo}{ext}"
                    ruta_archivo = tmp_dir / nombre_archivo_alt
                    contador += 1

            for intento in range(1, MAX_INTENTOS_DESCARGA + 1):
                try:
                    resp = session.get(
                        url,
                        stream=True,
                        timeout=30,
                        headers=REQUEST_HEADERS,
                        allow_redirects=True,
                    )

                    if resp.status_code == 200:
                        tam_header = resp.headers.get("content-length")
                        if tam_header:
                            try:
                                tam_esperado = int(tam_header)
                            except ValueError:
                                tam_esperado = None
                        else:
                            tam_esperado = None

                        with open(ruta_archivo, "wb") as f:
                            for chunk in resp.iter_content(chunk_size=CHUNK_SIZE):
                                if chunk:
                                    f.write(chunk)

                        tam_real = os.path.getsize(ruta_archivo)

                        if (tam_esperado is not None and tam_esperado == 0) or tam_real == 0:
                            try:
                                os.remove(ruta_archivo)
                            except OSError:
                                pass
                            raise Exception(
                                f"Archivo descargado con tamaño 0 (tam_esperado={tam_esperado}, tam_real={tam_real})."
                            )

                        descargado_ok = True
                        resultados.append(
                            {
                                "url": url,
                                "nombre_archivo": ruta_archivo.name,
                                "ruta_archivo": str(ruta_archivo),
                                "status": "OK",
                            }
                        )
                        break
                    else:
                        raise Exception(f"Código HTTP: {resp.status_code}")

                except Exception as e:
                    ultimo_error = str(e)
                    logger.warning(f"Error en intento {intento} para {nombre_archivo}: {e}")

                    if ruta_archivo and ruta_archivo.exists():
                        try:
                            os.remove(ruta_archivo)
                        except OSError:
                            pass

                    if intento < MAX_INTENTOS_DESCARGA:
                        espera = min(60, 2 ** intento)
                        logger.info(f"Reintentando en {espera} segundos para {url}...")
                        time.sleep(espera)

            if not descargado_ok:
                fallidos.append(
                    {
                        "url": url,
                        "nombre_archivo": nombre_archivo,
                        "error": ultimo_error or "Error desconocido",
                    }
                )

        except Exception as e:
            logger.error(f"Error al procesar {url}: {e}")
            fallidos.append(
                {
                    "url": url,
                    "nombre_archivo": nombre_archivo,
                    "error": str(e),
                }
            )

        # Progreso Streamlit
        elapsed = time.time() - start_time
        processed = idx
        pct = processed / total
        speed = processed / elapsed if elapsed > 0 else 0.0
        eta = (total - processed) / speed if speed > 0 else 0.0

        render_progress_bar_ui_task(progress_bar, pct)
        progress_text.markdown(
            f"Descargando Archivos {pct*100:.1f}% | {processed}/{total} "
            f"[{_format_hms(elapsed)}<{_format_hms(eta)}, {speed:.2f} archivo/s]"
        )


    # CSV de fallidos dentro del ZIP, si hay
    csv_fallidos_path = None
    if fallidos:
        csv_fallidos_path = tmp_dir / "descargas_fallidas.csv"
        pd.DataFrame(fallidos).to_csv(csv_fallidos_path, index=False, encoding="utf-8-sig")

    if not resultados and not csv_fallidos_path:
        return resultados, fallidos, None

    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for r in resultados:
            ruta = r.get("ruta_archivo")
            if ruta and os.path.exists(ruta):
                zf.write(ruta, arcname=Path(ruta).name)
        if csv_fallidos_path and csv_fallidos_path.exists():
            zf.write(csv_fallidos_path, arcname=csv_fallidos_path.name)

    zip_buffer.seek(0)
    return resultados, fallidos, zip_buffer.getvalue()


# ------------------------------------------------------
# Wrapper para PDFs extraídos de ZIP (simula UploadedFile)
# ------------------------------------------------------
class InMemoryUploadedPDF:
    """
    Pequeño wrapper para tratar PDFs extraídos de un ZIP o descargados en lote
    como si fueran `UploadedFile` de Streamlit.

    Además guarda, opcionalmente, la URL origen (`source_url`) del documento
    para poder enlazar luego con el Excel de contenidos (name, link_class).
    """

    def __init__(self, name: str, data: bytes, source_url: Optional[str] = None):
        self.name = name
        self._data = data
        self.source_url = source_url

    def getbuffer(self):
        return self._data


# ------------------------------------------------------
# Wrapper para DOCX extraídos de ZIP (simula UploadedFile)
# ------------------------------------------------------
class InMemoryUploadedDOCX:
    """
    Wrapper para tratar DOCX generados desde PDF o extraídos de ZIP como si fueran
    `UploadedFile` de Streamlit.

    También guarda (opcionalmente) la `source_url` del PDF original, cuando proviene
    de la Descarga Masiva.
    """

    def __init__(self, name: str, data: bytes, source_url: Optional[str] = None):
        self.name = name
        self._data = data
        self.source_url = source_url

    def getbuffer(self):
        return self._data

class InMemoryUploadedPPTX:
    """
    Wrapper en memoria para tratar PPTX (directos o desde ZIP / Descarga Masiva)
    de forma similar a los UploadedFile de Streamlit.
    """

    def __init__(self, name: str, data: bytes, source_url: Optional[str] = None):
        self.name = name
        self._data = data
        self.source_url = source_url

    def getbuffer(self):
        return self._data


# ======================================================
# HELPERS PDF to Word Transform (UI)
# ======================================================

def _build_pdf_file_table(uploaded_pdfs: List["st.runtime.uploaded_file_manager.UploadedFile"]) -> pd.DataFrame:
    """Resumen (nombre, tamaño, páginas) de PDFs subidos."""
    rows = []
    for f in uploaded_pdfs:
        data = f.getbuffer()
        size_mb = len(data) / (1024 * 1024)
        try:
            if fitz is not None:
                with fitz.open(stream=bytes(data), filetype="pdf") as doc:
                    pages = len(doc)
            else:
                pages = "?"
        except Exception:
            pages = "?"
        rows.append({"Nombre": f.name, "Tamaño_MB": round(size_mb, 2), "Páginas": pages})
    return pd.DataFrame(rows)


def _build_word_file_table(uploaded_docs: List["st.runtime.uploaded_file_manager.UploadedFile"]) -> pd.DataFrame:
    """Resumen (nombre, tamaño) de DOCX subidos."""
    rows = []
    for f in uploaded_docs:
        data = f.getbuffer()
        size_mb = len(data) / (1024 * 1024)
        rows.append(
            {
                "Nombre": f.name,
                "Tamaño_MB": round(size_mb, 2),
            }
        )
    return pd.DataFrame(rows)

def _build_pptx_file_table(uploaded_docs: List["st.runtime.uploaded_file_manager.UploadedFile"]) -> pd.DataFrame:
    """Resumen (nombre, tamaño) de PPTX subidos."""
    rows = []
    for f in uploaded_docs:
        data = f.getbuffer()
        size_mb = len(data) / (1024 * 1024)
        rows.append(
            {
                "Nombre": f.name,
                "Tamaño_MB": round(size_mb, 2),
            }
        )
    return pd.DataFrame(rows)


def _is_page_break_paragraph(para) -> bool:
    """
    Detecta si un párrafo es solamente un salto de página (add_page_break de python-docx).
    """
    WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    BR_TAG = f"{{{WORD_NS}}}br"
    TYPE_ATTR = f"{{{WORD_NS}}}type"

    try:
        for run in para.runs:
            r = run._r  # elemento XML <w:r>
            for br in r.iter():
                if br.tag == BR_TAG:
                    br_type = br.get(TYPE_ATTR)
                    # Consideramos salto de página el tipo "page" o sin tipo (por seguridad)
                    if br_type is None or br_type == "page":
                        return True
    except Exception:
        return False

    return False


def _iter_paragraphs_with_page(doc) -> List[Tuple[int, "docx.text.paragraph.Paragraph"]]:
    """
    Recorre TODOS los párrafos del cuerpo del documento (incluye tablas),
    devolviendo (número_de_página, párrafo).

    La numeración de página se estima:
    - Página inicial = 1
    - Cada párrafo que es sólo un salto de página (add_page_break) incrementa el contador
      y no se considera contenido.
    """
    from docx.text.paragraph import Paragraph  # import local para evitar problemas si falta python-docx

    WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    P_TAG = f"{{{WORD_NS}}}p"

    current_page = 1
    body = doc.element.body

    for elem in body.iter():
        if elem.tag != P_TAG:
            continue

        para = Paragraph(elem, doc)

        # Si es un párrafo que sólo representa un salto de página, avanzamos página y lo saltamos
        if _is_page_break_paragraph(para) and not (para.text or "").strip():
            current_page += 1
            continue

        yield current_page, para


def _extract_urls_from_text(text: str) -> List[str]:
    """
    Extrae URLs en texto plano, corrigiendo enlaces cortados por guion + espacio
    o guion + salto de línea (ej. '...spij-ext- web/#/...' o '...cgi- bin/...').

    Soporta prefijos:
      - http://, https://
      - www.
      - mailto:
      - tel:
    """
    s = _strip_invisible(text or "")
    if not s:
        return []

    # Prefijos reconocidos como inicio de URL
    prefixes = ("http://", "https://", "www.", "mailto:", "tel:")

    # Conjunto de caracteres válidos dentro de una URL (sin espacios)
    allowed = set("abcdefghijklmnopqrstuvwxyz"
                  "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
                  "0123456789"
                  "-._~:/?#[]@!$&'()*+,;=%")

    urls: List[str] = []
    i, n = 0, len(s)

    while i < n:
        lower = s[i:].lower()
        prefix = next((p for p in prefixes if lower.startswith(p)), None)
        if not prefix:
            i += 1
            continue

        # Hemos encontrado el inicio de una URL
        start = i
        j = i + len(prefix)
        url_chars = list(s[start:j])
        last_char = url_chars[-1] if url_chars else ""

        while j < n:
            c = s[j]

            # Caracter normal permitido en una URL
            if c in allowed:
                url_chars.append(c)
                last_char = c
                j += 1
                continue

            # Espacios (incluye saltos de línea internos de Word)
            if c.isspace():
                # Miramos más adelante saltando todos los espacios
                k = j + 1
                while k < n and s[k].isspace():
                    k += 1

                # Caso especial: URL cortada por guion de final de línea,
                # p.ej. '...spij-ext-' + salto + 'web/#/...'  o 'cgi-' + salto + 'bin/...'
                if k < n and last_char == "-" and s[k] in allowed:
                    # No añadimos el espacio, continuamos la URL directamente
                    j = k
                    continue

                # Si no es el caso anterior, terminamos la URL en el último carácter válido
                break

            # Puntuación habitual que cierra una URL
            if c in ".,;:!?)[]{}\"'":
                break

            # Otro carácter raro -> consideramos fin del enlace
            break

        # Construimos la URL y limpiamos puntuación de cierre
        url = "".join(url_chars).rstrip(".,;:!?)[]{}\"'")
        if url:
            urls.append(url)

        # Continuamos después de lo que hemos analizado
        i = j

    return urls



def _extract_urls_from_paragraph_xml(para, doc_part) -> List[str]:
    """
    Extrae URLs desde la estructura XML del párrafo:
    - w:hyperlink con r:id → relationships (target_ref)
    - campos HYPERLINK en w:instrText
    """
    urls: List[str] = []

    WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    HYPERLINK_TAG = f"{{{WORD_NS}}}hyperlink"
    INSTR_TEXT_TAG = f"{{{WORD_NS}}}instrText"
    RID_ATTR = f"{{{REL_NS}}}id"

    p = para._p  # elemento XML <w:p>

    # 1) Hyperlinks explícitos (target_ref en relationships)
    try:
        for hlink in p.iter():
            if hlink.tag != HYPERLINK_TAG:
                continue
            r_id = hlink.get(RID_ATTR)
            if not r_id:
                continue
            rel = doc_part.rels.get(r_id)
            if not rel:
                continue
            target = str(getattr(rel, "target_ref", "") or "")
            if not target:
                continue
            target = _strip_invisible(target).strip()
            target = target.strip("()[]{}.,;:!?'\"")
            if target:
                urls.append(target)
    except Exception:
        pass

    # 2) Campos HYPERLINK (instrText)
    try:
        for instr in p.iter():
            if instr.tag != INSTR_TEXT_TAG:
                continue
            txt = instr.text or ""
            for u in _extract_urls_from_text(txt):
                urls.append(u)
    except Exception:
        pass

    return urls


def _extract_links_from_docx_bytes(docx_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extrae todos los links de un DOCX, devolviendo filas:
    {
        "Nombre del Archivo": <nombre>,
        "Página/Diapositiva": <número de página estimado>,
        "Links": <url>
    }
    """
    if Document is None:
        return []

    from io import BytesIO

    try:
        doc = Document(BytesIO(docx_bytes))
    except Exception as e:
        logger.error(f"Error abriendo DOCX {filename}: {e}")
        return []

    rows: List[Dict[str, Any]] = []
    seen: set[Tuple[int, str]] = set()

    for page_idx, para in _iter_paragraphs_with_page(doc):
        # 1) URLs por texto plano (visible en el párrafo)
        urls_text = _extract_urls_from_text(para.text or "")

        # 2) URLs por estructura XML (hyperlinks, campos HYPERLINK)
        urls_xml = _extract_urls_from_paragraph_xml(para, doc.part)

        all_urls = urls_text + urls_xml

        for url in all_urls:
            if not url:
                continue
            key = (page_idx, url)
            if key in seen:
                continue
            seen.add(key)
            rows.append(
                {
                    "Nombre del Archivo": filename,
                    "Página/Diapositiva": str(page_idx),
                    "Links": url,
                }
            )

    return rows


def _run_word_link_report_streamlit(
    uploaded_docs: List["st.runtime.uploaded_file_manager.UploadedFile"],
    *,
    progress_bar,
    status_text,
) -> Tuple[pd.DataFrame, List[Dict[str, Any]]]:
    """
    Recorre todos los DOCX subidos (directos o desde ZIP) y arma
    un DataFrame con columnas:
    - Nombre del Archivo
    - Página/Diapositiva
    - Links
    - source_url (opcional, cuando proviene de Descarga Masiva)

    Devuelve (df_links, lista_errores).
    """
    if not uploaded_docs:
        empty_df = pd.DataFrame(columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"])
        return empty_df, []

    all_rows: List[Dict[str, Any]] = []
    errores: List[Dict[str, Any]] = []

    total_files = len(uploaded_docs)
    start_time = datetime.now()

    for idx, up in enumerate(uploaded_docs, start=1):
        file_name = up.name
        source_url = getattr(up, "source_url", None)

        status_text.markdown(f"Analizando **{idx}/{total_files}** · `{file_name}`")
        progress_bar.progress((idx - 1) / total_files)

        try:
            data = up.getbuffer()
            rows = _extract_links_from_docx_bytes(bytes(data), file_name)

            # 🔹 Añadimos la URL origen (si existe) a cada fila de links
            for r in rows:
                r["source_url"] = source_url

            all_rows.extend(rows)
        except Exception as e:
            logger.error(f"Error procesando Word {file_name}: {e}")
            errores.append({"Archivo": file_name, "Error": str(e)})

        render_progress_bar_ui_task(progress_bar, idx / total_files)

    if all_rows:
        df = pd.DataFrame(all_rows)
        df = df.sort_values(
            ["Nombre del Archivo", "Página/Diapositiva", "Links"]
        ).reset_index(drop=True)
    else:
        df = pd.DataFrame(columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"])

    elapsed = (datetime.now() - start_time).total_seconds()
    logger.info("Análisis de DOCX completado en %.2fs", elapsed)

    return df, errores

def _extract_links_from_pptx_bytes(pptx_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extrae todos los links de un PPTX, devolviendo filas:
    {
        "Nombre del Archivo": <nombre>,
        "Página/Diapositiva": <número de diapositiva>,
        "Links": <url>
    }

    Se combinan:
    - hipervínculos de texto (runs con hyperlink.address)
    - acciones de clic de forma (shape.click_action.hyperlink.address)
    - URLs en texto plano detectadas por regex (_extract_urls_from_text)
    """
    if Presentation is None:
        logger.warning("No se puede procesar PPTX porque falta la librería `python-pptx`.")
        return []

    from io import BytesIO

    try:
        prs = Presentation(BytesIO(pptx_bytes))
    except Exception as e:
        logger.error(f"Error abriendo PPTX {filename}: {e}")
        return []

    rows: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        urls_en_diapositiva: set[str] = set()

        for shape in slide.shapes:
            # 1) Hipervínculos asociados a acciones de clic de la forma
            try:
                click_action = getattr(shape, "click_action", None)
                if click_action is not None:
                    hlink = getattr(click_action, "hyperlink", None)
                    if hlink is not None and getattr(hlink, "address", None):
                        addr = str(hlink.address).strip()
                        if addr:
                            urls_en_diapositiva.add(addr)
            except Exception:
                pass

            # 2) Texto del shape (hipervínculos de runs y URLs visibles)
            text_for_plain = ""
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_for_plain = shape.text or ""

                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            h = getattr(run, "hyperlink", None)
                            if h is not None and getattr(h, "address", None):
                                addr = str(h.address).strip()
                                if addr:
                                    urls_en_diapositiva.add(addr)
            except Exception:
                pass

            if text_for_plain:
                for u in _extract_urls_from_text(text_for_plain):
                    urls_en_diapositiva.add(u)

        for url in sorted(urls_en_diapositiva):
            rows.append(
                {
                    "Nombre del Archivo": filename,
                    "Página/Diapositiva": str(slide_index),
                    "Links": url,
                }
            )

    return rows


def _run_pptx_link_report_streamlit(
    uploaded_pptx: List["InMemoryUploadedPPTX"],
    *,
    progress_bar,
    status_text,
) -> Tuple[pd.DataFrame, List[Dict[str, Any]]]:
    """
    Recorre todos los PPTX subidos (directos o desde ZIP) y arma
    un DataFrame con columnas:
    - Nombre del Archivo
    - Página/Diapositiva
    - Links
    - source_url (opcional, cuando proviene de Descarga Masiva)
    """
    if not uploaded_pptx:
        empty_df = pd.DataFrame(
            columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"]
        )
        return empty_df, []

    all_rows: List[Dict[str, Any]] = []
    errores: List[Dict[str, Any]] = []

    total_files = len(uploaded_pptx)
    start_time = datetime.now()

    for idx, up in enumerate(uploaded_pptx, start=1):
        file_name = up.name
        source_url = getattr(up, "source_url", None)

        status_text.markdown(f"Analizando PPTX **{idx}/{total_files}** · `{file_name}`")
        render_progress_bar_ui_task(progress_bar, (idx - 1) / total_files)

        try:
            data = up.getbuffer()
            rows = _extract_links_from_pptx_bytes(bytes(data), file_name)

            for r in rows:
                r["source_url"] = source_url

            all_rows.extend(rows)
        except Exception as e:
            logger.error(f"Error procesando PPTX {file_name}: {e}")
            errores.append({"Archivo": file_name, "Error": str(e)})

        render_progress_bar_ui_task(progress_bar, idx / total_files)

    if all_rows:
        df = pd.DataFrame(all_rows)
        df = df.sort_values(
            ["Nombre del Archivo", "Página/Diapositiva", "Links"]
        ).reset_index(drop=True)
    else:
        df = pd.DataFrame(
            columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"]
        )

    elapsed = (datetime.now() - start_time).total_seconds()
    logger.info("Análisis de PPTX completado en %.2fs", elapsed)

    return df, errores

def _run_pdf_extraction_streamlit(
    uploaded_pdfs: List["st.runtime.uploaded_file_manager.UploadedFile"],
    usar_multihilo: bool,
    max_workers: int,
    progress_bar,
    status_text,
):
    """
    Procesa en lote los PDFs subidos y devuelve (resultados, errores, zip_bytes).

    Los DOCX dentro del ZIP final usan el nombre original del PDF
    (o del PDF dentro del ZIP de entrada), sin prefijos como '001_'
    ni el prefijo 'Descarga_Masiva_Documentos_YYYYMMDD_HHMMSS_'.

    Además, cuando el PDF proviene de la Descarga Masiva, se conserva
    la URL origen (`source_url`) para poder enlazar con el Excel de URLs.
    """
    if not uploaded_pdfs:
        return [], [], None

    tmp_dir = Path(tempfile.mkdtemp(prefix="utp_pdf_extr_"))
    processor = PDFBatchProcessor(max_workers=max_workers)
    options = {
        "usar_multihilo": usar_multihilo,
        "max_workers": max_workers,
        # 🔧 IMPORTANTE: no filtrar bibliografía para conservar todos los links
        "filtrar_bibliografia": False,
    }

    resultados: List[Dict[str, Any]] = []
    errores: List[Dict[str, Any]] = []
    total_files = len(uploaded_pdfs)
    start_time = datetime.now()

    for idx, up in enumerate(uploaded_pdfs, start=1):
        pdf_name = up.name  # nombre original "visible" del PDF
        # 🔹 Si proviene de la Descarga Masiva, tendrá la URL origen
        source_url = getattr(up, "source_url", None)

        status_text.markdown(f"Procesando **{idx}/{total_files}** · `{pdf_name}`")
        render_progress_bar_ui_task(progress_bar, (idx - 1) / total_files)

        # PDF temporal (puede seguir llevando el índice para evitar colisiones)
        safe_stem = Path(pdf_name).stem
        pdf_temp_path = tmp_dir / f"{idx:03d}_{safe_stem}.pdf"
        with open(pdf_temp_path, "wb") as fh:
            fh.write(up.getbuffer())

        # Procesar PDF → Word
        result = processor.process_single_pdf(str(pdf_temp_path), str(tmp_dir), options)

        # Guardamos el nombre original del PDF y la URL origen
        result["original_pdf_name"] = pdf_name
        result["source_url"] = source_url

        if result.get("status") == "success":
            resultados.append(result)
        elif result.get("status") == "error":
            errores.append(result)
        else:
            errores.append(result)

        render_progress_bar_ui_task(progress_bar, idx / total_files)

    elapsed = (datetime.now() - start_time).total_seconds()
    logger.info("Procesamiento PDFs completado en %.2fs", elapsed)

    if not resultados:
        return resultados, errores, None

    # --- Construir ZIP con nombres limpios ---
    zip_buffer = io.BytesIO()

    used_names: set[str] = set()

    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        for r in resultados:
            out_path = r.get("output")
            if not out_path or not os.path.exists(out_path):
                continue

            # Nombre "ideal" del DOCX en el ZIP basado en el PDF original
            orig_name = r.get("original_pdf_name")
            if orig_name:
                orig_stem = Path(orig_name).stem

                # 1) quitar prefijo del ZIP de Descarga Masiva, si existe
                clean_stem = DESCARGA_PREFIX_RE.sub("", orig_stem)

                # por seguridad nunca dejarlo vacío
                if not clean_stem:
                    clean_stem = orig_stem

                base_name = clean_stem
            else:
                # Fallback: usar el nombre actual del DOCX
                base_name = Path(out_path).stem

            # 2) construir nombre final del DOCX
            arcname = f"{base_name}.docx"

            # 3) evitar colisiones dentro del ZIP
            if arcname in used_names:
                counter = 1
                while True:
                    alt = f"{base_name}_{counter}.docx"
                    if alt not in used_names:
                        arcname = alt
                        break
                    counter += 1
            used_names.add(arcname)

            # Escribimos el DOCX en el ZIP con el nombre final
            zf.write(out_path, arcname)

    zip_buffer.seek(0)

    # --- Crear lista de DOCX "en memoria" para el siguiente paso (Word → Links) ---
    word_docs: List[InMemoryUploadedDOCX] = []
    for r in resultados:
        out_path = r.get("output")
        if not out_path:
            continue
        try:
            if os.path.exists(out_path):
                with open(out_path, "rb") as fh:
                    data = fh.read()

                orig_name = r.get("original_pdf_name") or Path(out_path).name
                display_stem = Path(orig_name).stem
                # quitar prefijo Descarga_Masiva_Documentos_YYYYMMDD_HHMMSS_ si aplica
                display_stem = DESCARGA_PREFIX_RE.sub("", display_stem)
                if not display_stem:
                    display_stem = Path(out_path).stem

                display_name = f"{display_stem}.docx"

                # 🔹 Pasamos también la URL origen al DOCX "virtual"
                word_docs.append(
                    InMemoryUploadedDOCX(
                        display_name,
                        data,
                        source_url=r.get("source_url"),
                    )
                )
        except Exception as e:
            logger.error(f"Error leyendo DOCX generado '{out_path}': {e}")

    # Guardamos los DOCX en sesión para el módulo Word → Links
    st.session_state["pipeline_word_docs"] = word_docs

    return resultados, errores, zip_buffer.getvalue()


# ======================================================
# PÁGINAS / MÓDULOS
# ======================================================

def page_home():
    render_hero(
        title=APP_TITLE,
        subtitle="Plataforma de validación inteligente de Links rotos con reporte único de estado en Excel.",
        icon="🔗",
    )
    ui_card_open()
    st.markdown("### 🏠 Home")
    st.info(
        "Flujo recomendado: **Reporte Link (Excel) → Report Broken Link → Descargar Excel Status**"
    )
    ui_card_close()

def page_report_broken_unificado():
    """
    Pantalla unificada 'Report Broken Link':
    1) Bulk Download desde Excel
    2) PDF → Word
    3) Word → Links
    4) Links → Validación + Excel Status
    """

    # ======================================================
    # 1. Bulk Document (PDF, WORD and PPT) Download
    # ======================================================
    render_hero(
        "Bulk Document (PDF, WORD and PPT) Download",
        "Descarga múltiple de documentos PDF, Word y PPT desde un archivo Excel de URLs.",
        "⬇️",
    )

    # Botón Reiniciar justo debajo del primer bloque
    st.markdown('<div class="hero-reset-anchor"></div>', unsafe_allow_html=True)
    if st.button("Reiniciar", key="btn_reset_report_broken"):
        reset_report_broken_pipeline()
        try:
            st.rerun()
        except AttributeError:
            st.experimental_rerun()

    # ---------- TARJETA: Selección de Excel (Paso 1) ----------
    ui_card_open()
    step1_ph = st.empty()

    bulk_uploader_key = f"pipeline_bulk_excel_uploader_{st.session_state.get('pipeline_reset_token', 0)}"
    uploaded_excel = st.file_uploader(
        "Seleccione el archivo Excel que contiene las URLs de los documentos a descargar",
        type=["xlsx", "xls"],
        key=bulk_uploader_key,
    )

    file_ok = uploaded_excel is not None
    step1_ph.markdown(
        render_step_header_html(
            "1",
            "Seleccione el archivo Excel que contiene las URLs de los documentos a descargar",
            "ok" if file_ok else "warn",
        ),
        unsafe_allow_html=True,
    )

    bulk_urls_archivos: List[str] = []
    df_in_bulk: Optional[pd.DataFrame] = None

    st.session_state["bulk_has_valid_urls"] = False
    st.session_state["bulk_urls_archivos"] = None
    st.session_state["bulk_excel_df"] = None 

    if file_ok:
        try:
            df_in_bulk = _read_excel_safe(uploaded_excel)
        except Exception as e:
            st.error(str(e))
            
        else:
            st.session_state["bulk_excel_df"] = df_in_bulk
            
            if "url" not in df_in_bulk.columns:
                st.error("El Excel no contiene la columna requerida: **url**.")
                st.caption(
                    f"Columnas detectadas: {', '.join(map(str, df_in_bulk.columns.tolist()))}"
                )
            else:
                df_urls = df_in_bulk["url"].dropna().astype(str).str.strip()
                bulk_urls_archivos = [
                    u for u in df_urls if u.lower().endswith(DESC_EXT_PERMITIDAS)
                ]
                total_urls = len(df_urls)
                total_permitidas = len(bulk_urls_archivos)

                if total_permitidas == 0:
                    st.warning(
                        "No se encontraron URLs que terminen en .ppt, .pptx, .pdf, .doc o .docx."
                    )
                else:
                    st.session_state["bulk_has_valid_urls"] = True
                    st.session_state["bulk_urls_archivos"] = bulk_urls_archivos

                    # Firma para detectar cambio de Excel
                    try:
                        excel_bytes = uploaded_excel.getbuffer()
                        bulk_signature = (uploaded_excel.name, len(excel_bytes))
                    except Exception:
                        bulk_signature = (uploaded_excel.name, 0)

                    prev_bulk_sig = st.session_state.get("pipeline_bulk_signature")
                    if prev_bulk_sig != bulk_signature:
                        st.session_state["pipeline_bulk_signature"] = bulk_signature
                        st.session_state["pipeline_bulk_done"] = False

                        # Reset fases posteriores
                        st.session_state["pipeline_pdf_signature"] = None
                        st.session_state["pipeline_pdf_done"] = False
                        st.session_state["pipeline_pdf_results"] = None
                        st.session_state["pipeline_pdf_errors"] = None

                        st.session_state["pipeline_word_docs"] = None
                        st.session_state["pipeline_word_done"] = False
                        st.session_state["pipeline_df_links"] = None
                        st.session_state["pipeline_word_errors"] = None

                        st.session_state["pipeline_status_done"] = False
                        st.session_state["status_result_df"] = None
                        st.session_state["status_export_df"] = None
                        st.session_state["status_invalid_df"] = None

    else:
        st.caption(
            "Carga un archivo Excel para continuar con la descarga masiva."
        )

    ui_card_close()

    # ---------- TARJETA: Procesar Descarga Masiva (Paso 2) ----------
    ui_card_open()
    render_simple_step_header("2", "Procesar Descarga Masiva")

    if not _requests_available_or_warn():
        ui_card_close()
        return

    progress_bar_bulk = st.empty()
    progress_text_bulk = st.empty()

    urls_archivos_state = st.session_state.get("bulk_urls_archivos") or []
    auto_trigger_bulk = bool(urls_archivos_state) and not st.session_state.get(
        "pipeline_bulk_done", False
    )
    manual_click_bulk = st.button(
        "🚀 Procesar Descarga Masiva",
        type="primary",
        key="pipeline_btn_bulk_process",
    )

    if urls_archivos_state and (auto_trigger_bulk or manual_click_bulk):
        try:
            render_progress_bar_ui_task(progress_bar_bulk, 0.0)
            progress_text_bulk.markdown("Preparando descarga masiva...")


            with st.spinner("Descargando archivos..."):
                resultados, fallidos, zip_bytes = _run_descarga_masiva_streamlit(
                    urls_archivos_state,
                    progress_bar=progress_bar_bulk,
                    progress_text=progress_text_bulk,
                )

            st.session_state["descarga_resultados"] = resultados
            st.session_state["descarga_fallidos"] = fallidos
            st.session_state["descarga_zip_bytes"] = zip_bytes
            st.session_state["pipeline_bulk_done"] = True

            progress_bar_bulk.empty()
            progress_text_bulk.markdown("✅ Descarga masiva completada.")
        except Exception as e:
            progress_bar_bulk.empty()
            progress_text_bulk.empty()
            st.error(f"Ocurrió un error durante la descarga masiva: {e}")

    # Pequeño resumen (sin expander de detalle)
    resultados_ready = st.session_state.get("descarga_resultados") or []
    fallidos_ready = st.session_state.get("descarga_fallidos") or []

    # Ya no mostramos métricas de Total archivos / OK / Con error para este paso.
    # Solo mostramos el mensaje de ayuda cuando aún no hay URLs para procesar.
    if not (resultados_ready or fallidos_ready):
        if not urls_archivos_state:
            st.caption(
                "Primero carga un Excel válido con URLs para poder procesar la descarga."
            )


    ui_card_close()

    # ---------- TARJETA: Descargar ZIP (Paso 3) ----------
    ui_card_open()
    render_simple_step_header("3", "Descargar todos los archivos (PDF, Word, PPT) (ZIP)")

    zip_bytes_ready = st.session_state.get("descarga_zip_bytes")
    resultados_ready = st.session_state.get("descarga_resultados") or []
    fallidos_ready = st.session_state.get("descarga_fallidos") or []

    if not zip_bytes_ready or (not resultados_ready and not fallidos_ready):
        st.warning(
            "Primero ejecuta el paso 2 para generar las descargas."
        )
    else:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        zip_name = f"Descarga_Masiva_Documentos_{ts}.zip"
        st.download_button(
            "⬇️ Descargar todos los archivos (ZIP)",
            data=zip_bytes_ready,
            file_name=zip_name,
            mime="application/zip",
        )

    ui_card_close()

     # ======================================================
    # 2. PDF, WORD and PPT to Word Transformation (ZIP)
    # ======================================================
    render_hero(
        "PDF, WORD and PPT to Word Transformation (ZIP)",
        "Convierte múltiples PDFs, Word y PPT a Word (texto) sin filtrar bibliografía y reordenando el texto.",
        "🧲",
    )

    # Validar dependencias básicas para este módulo
    if fitz is None or Document is None:
        ui_card_open()
        st.error(
            "Faltan dependencias para este módulo.\n\n"
            "- Instala `pymupdf` (fitz)\n"
            "- Instala `python-docx`\n\n"
            "Luego vuelve a desplegar la aplicación."
        )
        ui_card_close()
        return

    ui_card_open()
    step_pdf1 = st.empty()

    pdf_uploader_key = f"pipeline_pdf_uploader_ultra_{st.session_state.get('pipeline_reset_token', 0)}"
    uploaded_files = st.file_uploader(
        "Selecciona uno o más archivos PDF, Word, PPT o ZIP (ZIP con PDFs/Word/PPT en su interior)",
        type=["pdf", "docx", "pptx", "doc", "ppt", "zip"],
        accept_multiple_files=True,
        key=pdf_uploader_key,
        help="Si ya ejecutaste la Descarga Masiva, aquí llegarán automáticamente los documentos.",
    )

    all_pdfs: List[Any] = []
    all_word_docs: List[InMemoryUploadedDOCX] = []
    all_pptx_docs: List["InMemoryUploadedPPTX"] = []
    unsupported_office: List[Dict[str, str]] = []

    # Documentos provenientes de Descarga Masiva (si existen)
    resultados_desc = st.session_state.get("descarga_resultados") or []
    for r in resultados_desc:
        ruta = r.get("ruta_archivo")
        if not ruta:
            continue
        ruta_str = str(ruta)
        ext = Path(ruta_str).suffix.lower()

        try:
            with open(ruta_str, "rb") as fh:
                data = fh.read()
        except Exception as e:
            logger.warning(f"No se pudo leer archivo descargado '{ruta_str}': {e}")
            continue

        source_url = r.get("url")
        file_name = Path(ruta_str).name

        if ext == ".pdf":
            all_pdfs.append(
                InMemoryUploadedPDF(
                    file_name,
                    data,
                    source_url=source_url,
                )
            )
        elif ext == ".docx":
            all_word_docs.append(
                InMemoryUploadedDOCX(
                    file_name,
                    data,
                    source_url=source_url,
                )
            )
        elif ext == ".pptx":
            if Presentation is not None:
                all_pptx_docs.append(
                    InMemoryUploadedPPTX(
                        file_name,
                        data,
                        source_url=source_url,
                    )
                )
            else:
                unsupported_office.append(
                    {
                        "Archivo": file_name,
                        "Motivo": "python-pptx no está instalado; no se pueden procesar PPTX descargados.",
                    }
                )
        elif ext in (".doc", ".ppt"):
            unsupported_office.append(
                {
                    "Archivo": file_name,
                    "Motivo": f"Formato {ext} no soportado. Convierte a .docx o .pptx para poder analizar los links.",
                }
            )

    # Documentos / ZIP subidos manualmente
    if uploaded_files:
        for f in uploaded_files:
            fname_lower = f.name.lower()
            ext = Path(fname_lower).suffix.lower()

            # PDF directos
            if ext == ".pdf":
                all_pdfs.append(
                    InMemoryUploadedPDF(
                        f.name,
                        f.getbuffer(),
                        source_url=f.name,
                    )
                )

            # Word (DOCX) directos
            elif ext == ".docx":
                all_word_docs.append(
                    InMemoryUploadedDOCX(
                        f.name,
                        f.getbuffer(),
                        source_url=f.name,
                    )
                )

            # PPTX directos
            elif ext == ".pptx":
                if Presentation is not None:
                    all_pptx_docs.append(
                        InMemoryUploadedPPTX(
                            f.name,
                            f.getbuffer(),
                            source_url=f.name,
                        )
                    )
                else:
                    unsupported_office.append(
                        {
                            "Archivo": f.name,
                            "Motivo": "python-pptx no está instalado; no se pueden procesar PPTX subidos manualmente.",
                        }
                    )

            # Formatos antiguos no soportados
            elif ext in (".doc", ".ppt"):
                unsupported_office.append(
                    {
                        "Archivo": f.name,
                        "Motivo": f"Formato {ext} no soportado. Convierte a .docx o .pptx antes de cargarlo.",
                    }
                )

            # ZIP con mezcla de documentos
            elif ext == ".zip":
                try:
                    zdata = io.BytesIO(f.getbuffer())
                    with zipfile.ZipFile(zdata, "r") as zf:
                        for info in zf.infolist():
                            if info.is_dir():
                                continue

                            inner_ext = Path(info.filename).suffix.lower()
                            inner_name = Path(info.filename).name
                            file_bytes = zf.read(info)

                            if inner_ext == ".pdf":
                                all_pdfs.append(
                                    InMemoryUploadedPDF(
                                        inner_name,
                                        file_bytes,
                                        source_url=f.name,
                                    )
                                )
                            elif inner_ext == ".docx":
                                all_word_docs.append(
                                    InMemoryUploadedDOCX(
                                        inner_name,
                                        file_bytes,
                                        source_url=f.name,
                                    )
                                )
                            elif inner_ext == ".pptx":
                                if Presentation is not None:
                                    all_pptx_docs.append(
                                        InMemoryUploadedPPTX(
                                            inner_name,
                                            file_bytes,
                                            source_url=f.name,
                                        )
                                    )
                                else:
                                    unsupported_office.append(
                                        {
                                            "Archivo": inner_name,
                                            "Motivo": "python-pptx no está instalado; no se pueden procesar PPTX dentro del ZIP.",
                                        }
                                    )
                            elif inner_ext in (".doc", ".ppt"):
                                unsupported_office.append(
                                    {
                                        "Archivo": inner_name,
                                        "Motivo": f"Formato {inner_ext} no soportado. Convierte a .docx o .pptx para poder analizar los links.",
                                    }
                                )
                except Exception as e:
                    st.warning(f"No se pudo leer el ZIP `{f.name}`: {e}")

    has_docs = (len(all_pdfs) + len(all_word_docs) + len(all_pptx_docs)) > 0

    step_pdf1.markdown(
        render_step_header_html(
            "4",
            "Agregar documentos (PDF, Word, PPT) directos o desde ZIP",
            "ok" if has_docs else "warn",
        ),
        unsafe_allow_html=True,
    )

    if not has_docs:
        st.caption(
            "Agrega documentos manualmente o ejecuta primero la Descarga Masiva para que lleguen aquí automáticamente."
        )
        ui_card_close()
        return

    # Firma de documentos para controlar re-ejecuciones
    all_inputs_for_signature: List[Any] = []
    all_inputs_for_signature.extend(all_pdfs)
    all_inputs_for_signature.extend(all_word_docs)
    all_inputs_for_signature.extend(all_pptx_docs)

    try:
        signature = sorted(
            (getattr(f, "name", str(idx)), len(f.getbuffer()))
            for idx, f in enumerate(all_inputs_for_signature)
        )
    except Exception:
        signature = sorted(
            (getattr(f, "name", str(idx)), 0) for idx, f in enumerate(all_inputs_for_signature)
        )

    prev_sig = st.session_state.get("pipeline_pdf_signature")
    if prev_sig != signature:
        st.session_state["pipeline_pdf_signature"] = signature
        st.session_state["pipeline_pdf_done"] = False
        st.session_state["pipeline_pdf_results"] = None
        st.session_state["pipeline_pdf_errors"] = None

        st.session_state["pipeline_word_docs"] = None
        st.session_state["pipeline_ppt_docs"] = None
        st.session_state["pipeline_word_done"] = False
        st.session_state["pipeline_df_links"] = None
        st.session_state["pipeline_word_errors"] = None

        st.session_state["pipeline_status_done"] = False
        st.session_state["status_result_df"] = None
        st.session_state["status_export_df"] = None
        st.session_state["status_invalid_df"] = None

    # Listado de archivos seleccionados
    if all_pdfs:
        with st.expander("Archivos PDF seleccionados", expanded=False):
            df_files = _build_pdf_file_table(all_pdfs)
            st.dataframe(df_files, use_container_width=True, height=260)

    if all_word_docs:
        with st.expander("Archivos Word (DOCX) seleccionados", expanded=False):
            df_files_word = _build_word_file_table(all_word_docs)
            st.dataframe(df_files_word, use_container_width=True, height=260)

    if all_pptx_docs:
        with st.expander("Archivos PPTX seleccionados", expanded=False):
            df_files_ppt = _build_pptx_file_table(all_pptx_docs)
            st.dataframe(df_files_ppt, use_container_width=True, height=260)

    if unsupported_office:
        with st.expander("⚠️ Archivos Office no soportados", expanded=False):
            df_unsup = pd.DataFrame(unsupported_office)
            st.dataframe(df_unsup, use_container_width=True, height=260)

    # Opciones de procesamiento (aplican solo a PDFs)
    with st.expander("Opciones de procesamiento de PDFs", expanded=False):
        col1, col2 = st.columns(2)
        with col1:
            usar_multihilo = st.toggle(
                "Usar procesamiento paralelo por páginas (solo PDFs)",
                value=True,
                help="Recomendado cuando los PDFs tienen muchas páginas.",
                key="extr_usar_multihilo",
            )
        with col2:
            max_workers = st.number_input(
                "Número máximo de workers para PDFs",
                min_value=1,
                max_value=16,
                value=4,
                step=1,
                key="extr_max_workers",
            )

    # 4. Procesar todos los documentos (Paso 5)
    render_simple_step_header("5", "Procesar todos los documentos (PDF, Word, PPT)")

    progress_bar_pdf = st.empty()
    status_text_pdf = st.empty()

    auto_trigger_pdf = not st.session_state.get("pipeline_pdf_done", False)
    manual_click_pdf = st.button(
        "🚀 Procesar todos los documentos",
        type="primary",
        key="pipeline_btn_pdf_process",
    )

    if auto_trigger_pdf or manual_click_pdf:
        try:
            render_progress_bar_ui_task(progress_bar_pdf, 0.0)
            status_text_pdf.markdown("Iniciando procesamiento de documentos...")


            resultados_pdf: List[Dict[str, Any]] = []
            errores_pdf: List[Dict[str, Any]] = []
            zip_bytes_pdf: Optional[bytes] = None

            # 1) PDFs → DOCX (texto)
            if all_pdfs:
                with st.spinner("Extrayendo texto y generando archivos Word desde PDFs..."):
                    resultados_pdf, errores_pdf, zip_bytes_pdf = _run_pdf_extraction_streamlit(
                        all_pdfs,
                        usar_multihilo=bool(usar_multihilo),
                        max_workers=int(max_workers),
                        progress_bar=progress_bar_pdf,
                        status_text=status_text_pdf,
                    )
            else:
                status_text_pdf.markdown("No hay PDFs para procesar; se usarán solo Word/PPT existentes.")

            # 2) DOCX generados desde PDFs + DOCX originales
            word_docs_from_pdf = st.session_state.get("pipeline_word_docs") or []
            if not isinstance(word_docs_from_pdf, list):
                word_docs_from_pdf = []

            combined_word_docs: List[InMemoryUploadedDOCX] = []
            combined_word_docs.extend(word_docs_from_pdf)
            combined_word_docs.extend(all_word_docs)

            st.session_state["pipeline_word_docs"] = combined_word_docs
            st.session_state["pipeline_ppt_docs"] = all_pptx_docs

            # Guardar resultados en sesión
            st.session_state["pipeline_pdf_results"] = resultados_pdf
            st.session_state["pipeline_pdf_errors"] = errores_pdf
            st.session_state["extraccion_resultados"] = resultados_pdf
            st.session_state["extraccion_errores"] = errores_pdf
            st.session_state["extraccion_zip_bytes"] = zip_bytes_pdf
            st.session_state["pipeline_pdf_done"] = True
            st.session_state["pipeline_word_inputs_count"] = len(all_word_docs)
            st.session_state["pipeline_ppt_inputs_count"] = len(all_pptx_docs)

            progress_bar_pdf.empty()
            status_text_pdf.markdown("✅ Procesamiento de documentos completado.")
        except Exception as e:
            progress_bar_pdf.empty()
            status_text_pdf.empty()
            st.error(f"Ocurrió un error durante el procesamiento de documentos: {e}")
    else:
        resultados_pdf = st.session_state.get("pipeline_pdf_results") or []
        errores_pdf = st.session_state.get("pipeline_pdf_errors") or []
        word_count = st.session_state.get("pipeline_word_inputs_count", 0)
        ppt_count = st.session_state.get("pipeline_ppt_inputs_count", 0)

        if resultados_pdf or errores_pdf or word_count or ppt_count:
            total_ok = len(resultados_pdf)
            total_err = len(errores_pdf)
            total_pdf_files = total_ok + total_err
            total_pag = sum(
                r.get("stats", {}).get("paginas_procesadas", 0)
                for r in resultados_pdf
                if r.get("stats")
            )

            m1, m2, m3, m4, m5 = st.columns(5)
            m1.metric("PDF procesados", total_pdf_files)
            m2.metric("PDF OK", total_ok)
            m3.metric("PDF con error", total_err)
            m4.metric("Word (DOCX) detectados", word_count)
            m5.metric("PPTX detectados", ppt_count)

            if total_pag:
                st.caption(f"Páginas totales procesadas en PDFs: **{total_pag}**")

    ui_card_close()


    # ======================================================
    # 3. Report Word & PPT Link (Word / PPT → Links)
    # ======================================================
    render_hero(
        "Report Word & PPT Link (Excel)",
        "Genera un Excel con todos los links detectados en documentos Word y PPT.",
        "📊",
    )

    ui_card_open()

    docs = st.session_state.get("pipeline_word_docs") or []
    ppt_docs = st.session_state.get("pipeline_ppt_docs") or []
    total_docs_inputs = len(docs) + len(ppt_docs)

    step_word1 = st.empty()
    step_word1.markdown(
        render_step_header_html(
            "6",
            "Agregar documentos Word's y PPT's (directos o desde ZIP)",
            "ok" if total_docs_inputs > 0 else "warn",
        ),
        unsafe_allow_html=True,
    )

    if total_docs_inputs == 0:
        st.caption(
            "Los documentos Word se generan automáticamente a partir de los PDFs del paso anterior; "
            "los Word/PPT que cargues o descargues también se incorporan aquí."
        )
        ui_card_close()
    else:
        if docs:
            with st.expander("Documentos Word (DOCX) a analizar", expanded=False):
                df_files = _build_word_file_table(docs)
                st.dataframe(df_files, use_container_width=True, height=260)

        if ppt_docs:
            with st.expander("Presentaciones PPTX a analizar", expanded=False):
                df_files_ppt = _build_pptx_file_table(ppt_docs)
                st.dataframe(df_files_ppt, use_container_width=True, height=260)

        # Procesar Word & PPT (Paso 7)
        render_simple_step_header("7", "Procesar todos los documentos Word's y PPT's")

        progress_bar_word = st.empty()
        status_text_word = st.empty()

        auto_trigger_word = not st.session_state.get("pipeline_word_done", False)
        manual_click_word = st.button(
            "🚀 Procesar todos los documentos Word's y PPT's",
            type="primary",
            key="pipeline_btn_word_process",
        )

        if auto_trigger_word or manual_click_word:
            try:
                render_progress_bar_ui_task(progress_bar_word, 0.0)
                status_text_word.markdown("Iniciando análisis de documentos Word y PPT...")

                df_links_word = pd.DataFrame(
                    columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"]
                )
                errores_word: List[Dict[str, Any]] = []

                if docs:
                    with st.spinner("Buscando enlaces dentro de los documentos Word..."):
                        df_links_word, errores_word = _run_word_link_report_streamlit(
                            docs,
                            progress_bar=progress_bar_word,
                            status_text=status_text_word,
                        )

                df_links_ppt = pd.DataFrame(
                    columns=["Nombre del Archivo", "Página/Diapositiva", "Links", "source_url"]
                )
                errores_ppt: List[Dict[str, Any]] = []

                if ppt_docs:
                    with st.spinner("Buscando enlaces dentro de las presentaciones PPTX..."):
                        df_links_ppt, errores_ppt = _run_pptx_link_report_streamlit(
                            ppt_docs,
                            progress_bar=progress_bar_word,
                            status_text=status_text_word,
                        )

                if not df_links_word.empty and not df_links_ppt.empty:
                    df_links = pd.concat([df_links_word, df_links_ppt], ignore_index=True, sort=False)
                elif not df_links_word.empty:
                    df_links = df_links_word
                else:
                    df_links = df_links_ppt

                errores = errores_word + errores_ppt

                st.session_state["pipeline_df_links"] = df_links
                st.session_state["reporte_links_df"] = df_links
                st.session_state["pipeline_word_errors"] = errores
                st.session_state["pipeline_word_done"] = True

                progress_bar_word.empty()
                status_text_word.markdown("✅ Análisis de documentos Word y PPT completado.")

                total_docs = len(docs) + len(ppt_docs)
                total_links = len(df_links)
                docs_con_links = (
                    df_links["Nombre del Archivo"].nunique()
                    if not df_links.empty
                    else 0
                )

                m1, m2, m3 = st.columns(3)
                m1.metric("Documentos analizados", total_docs)
                m2.metric("Documentos con links", docs_con_links)
                m3.metric("Links detectados", total_links)

                with st.expander("📊 Detalle de links detectados", expanded=False):
                    if not df_links.empty:
                        st.dataframe(df_links, use_container_width=True, height=320)
                    else:
                        st.write("No se detectaron links en los documentos analizados.")

            except Exception as e:
                progress_bar_word.empty()
                status_text_word.empty()
                st.error(f"Ocurrió un error durante el análisis de documentos: {e}")
        else:
            df_links = st.session_state.get("pipeline_df_links")
            if df_links is not None:
                total_docs = len(docs) + len(ppt_docs)
                total_links = len(df_links)
                docs_con_links = (
                    df_links["Nombre del Archivo"].nunique()
                    if not df_links.empty
                    else 0
                )

                m1, m2, m3 = st.columns(3)
                m1.metric("Documentos analizados", total_docs)
                m2.metric("Documentos con links", docs_con_links)
                m3.metric("Links detectados", total_links)

                with st.expander("📊 Detalle de links detectados", expanded=False):
                    if not df_links.empty:
                        st.dataframe(df_links, use_container_width=True, height=320)
                    else:
                        st.write("No se detectaron links en los documentos analizados.")
            else:
                st.info("Cuando termine el análisis se mostrará aquí el resumen de links.")

        ui_card_close()


    # ======================================================
    # 4. Report Broken Link (Excel → Status)
    # ======================================================
    render_hero(
        "Report Broken Link (Excel)",
        "Comprueba automáticamente si los Links de tu reporte Excel están activos o rotos.",
        "🧭",
    )

    # 4.1 / 4.2: Cargar reporte (desde el paso anterior) + Detalle de Links
    ui_card_open()
    step_excel1 = st.empty()
    df_links = st.session_state.get("pipeline_df_links")

    has_links = df_links is not None and not df_links.empty

    step_excel1.markdown(
        render_step_header_html(
            "8",
            "Procesar Reporte Link (desde documentos Word)",
            "ok" if has_links else "warn",
        ),
        unsafe_allow_html=True,
    )

    if not has_links:
        st.caption(
            "Primero procesa los documentos Word en el paso 7 para generar el reporte de links."
        )
        ui_card_close()
        return

    df_in = df_links.copy()
    df_in["Fila_Excel"] = range(2, 2 + len(df_in))

    default_scheme = "https"
    allow_mailto = False
    allow_tel = False
    allow_anchor = False

    links_with_rows, df_invalid = _normalize_links(
        df_in["Links"],
        allow_mailto=allow_mailto,
        allow_tel=allow_tel,
        allow_anchors_only=allow_anchor,
        default_scheme=default_scheme,
    )

    st.session_state.status_input_filename = "Reporte_Link_Automatizado"
    st.session_state.status_input_df = df_in
    st.session_state.status_links_list = links_with_rows
    st.session_state.status_result_df = None
    st.session_state.status_invalid_df = df_invalid
    st.session_state.status_export_df = None

    with st.expander("Detalle de Links", expanded=True):
        c1, c2, c3 = st.columns(3)
        c1.metric("Links válidos", f"{len(links_with_rows)}")
        c2.metric("Total filas", f"{len(df_in)}")
        c3.metric("Descartados", f"{len(df_invalid)}")

    if not df_invalid.empty:
        with st.expander("⚠️ Links descartados (Status = INVALIDO)", expanded=False):
            st.dataframe(df_invalid, use_container_width=True, height=220)

    ui_card_close()

    # 4.3 Procesar Reporte Link (Paso 8)
    ui_card_open()
    render_simple_step_header("9", "Procesar y Descargar Status Reporte (Excel)")

    if not _httpx_available_or_warn():
        ui_card_close()
        return

    if len(links_with_rows) == 0 and df_invalid.empty:
        st.warning("No hay links para procesar (ni válidos ni descartados).")
        ui_card_close()
        return

    concurrency_global = DEFAULT_CONCURRENCY_GLOBAL
    concurrency_per_host = DEFAULT_CONCURRENCY_PER_HOST
    timeout_s = DEFAULT_TIMEOUT_S
    retries = DEFAULT_RETRIES
    detect_soft_404 = True
    verify_ssl = True
    max_bytes = DEFAULT_MAX_BYTES
    range_bytes = DEFAULT_RANGE_BYTES

    progress_bar = st.empty()
    status_text = st.empty()

    def progress_cb(done: int, total: int, current_url: str, current_status: str):
        pct = done / max(1, total)
        render_progress_bar_ui_task(progress_bar, pct)
        show = current_url if len(current_url) <= 85 else ("…" + current_url[-82:])
        status_text.markdown(
            f"Validando **{done}/{total}** · `{show}` · **{current_status}**"
        )

    should_run = False
    if not st.session_state.get("pipeline_status_done", False):
        should_run = True

    if st.button("🚀 Iniciar validación", type="primary", key="pipeline_btn_status_process"):
        should_run = True
        st.session_state["pipeline_status_done"] = False

    df_out: Optional[pd.DataFrame] = None

    if should_run:
        try:
            render_progress_bar_ui_task(progress_bar, 0.0)
            status_text.markdown("Iniciando verificación con motor...")

            if len(links_with_rows) > 0:
                with st.spinner(
                    "Validando enlaces (scoring, dominios especiales, soft-404 mejorado)..."
                ):
                    results = run_async(
                        _run_link_check_ultra_v5(
                            links_with_rows,
                            timeout_s=float(timeout_s),
                            concurrency_global=int(concurrency_global),
                            concurrency_per_host=int(concurrency_per_host),
                            detect_soft_404=bool(detect_soft_404),
                            retries=int(retries),
                            verify_ssl=bool(verify_ssl),
                            max_bytes=int(max_bytes),
                            range_bytes=int(range_bytes),
                            progress_callback=progress_cb,
                        )
                    )
                df_out = pd.DataFrame(results)
            else:
                df_out = pd.DataFrame(
                    columns=[
                        "Link",
                        "Status",
                        "HTTP_Code",
                        "Detalle",
                        "Content_Type",
                        "Redirected",
                        "Timestamp",
                        "Final_URL",
                        "Redirect_Chain",
                        "Soft_404",
                        "Score",
                        "Fila_Excel",
                    ]
                )

            # --- 1) Construir filas para links inválidos (si los hay) ---
            df_invalid_ready = st.session_state.get("status_invalid_df")
            if isinstance(df_invalid_ready, pd.DataFrame) and not df_invalid_ready.empty:
                now_str_invalid = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                invalid_rows: List[Dict[str, Any]] = []

                for _, inv in df_invalid_ready.iterrows():
                    invalid_rows.append(
                        {
                            "Link": inv["Valor"],
                            "Status": "INVALIDO",
                            "HTTP_Code": None,
                            "Detalle": inv["Motivo"],
                            "Content_Type": "",
                            "Redirected": "No",
                            "Timestamp": now_str_invalid,
                            "Final_URL": inv["Valor"],
                            "Redirect_Chain": inv["Valor"],
                            "Soft_404": "No",
                            "Score": -100,
                            "Fila_Excel": inv["Fila_Excel"],
                        }
                    )

                df_invalid_status = pd.DataFrame(
                    invalid_rows,
                    columns=[
                        "Link",
                        "Status",
                        "HTTP_Code",
                        "Detalle",
                        "Content_Type",
                        "Redirected",
                        "Timestamp",
                        "Final_URL",
                        "Redirect_Chain",
                        "Soft_404",
                        "Score",
                        "Fila_Excel",
                    ],
                )

                df_out = pd.concat([df_out, df_invalid_status], ignore_index=True, sort=False)

            # --- 2) Añadir SIEMPRE metadatos del reporte de links (Word → Links) ---
            df_src = df_in.copy()
            meta_cols: List[str] = []
            if "Nombre del Archivo" in df_src.columns:
                meta_cols.append("Nombre del Archivo")
            if "Página/Diapositiva" in df_src.columns:
                meta_cols.append("Página/Diapositiva")
            # conservar también la URL origen del documento (PDF descargado)
            if "source_url" in df_src.columns:
                meta_cols.append("source_url")

            if meta_cols:
                df_meta = df_src[["Fila_Excel"] + meta_cols]
                df_out = df_out.merge(df_meta, on="Fila_Excel", how="left")

            # --- 3) Enriquecer con name y link_class desde el Excel de URLs (si existe) ---
            bulk_df = st.session_state.get("bulk_excel_df")
            if isinstance(bulk_df, pd.DataFrame) and "url" in bulk_df.columns:
                cols_to_add: List[str] = []
                if "name" in bulk_df.columns:
                    cols_to_add.append("name")
                if "link_class" in bulk_df.columns:
                    cols_to_add.append("link_class")

                if cols_to_add and "source_url" in df_out.columns:
                    df_bulk = bulk_df[["url"] + cols_to_add].copy()
                    df_out = df_out.merge(
                        df_bulk,
                        left_on="source_url",
                        right_on="url",
                        how="left",
                    )
                    # ya no necesitamos la columna 'url' del Excel origen
                    if "url" in df_out.columns:
                        df_out = df_out.drop(columns=["url"])

            # --- 4) Asegurar existencia de columnas name / link_class aunque no haya Excel ---
            for col_extra in ("name", "link_class"):
                if col_extra not in df_out.columns:
                    df_out[col_extra] = ""

            # --- 5) Rellenar 'name' con el nombre del ZIP/PDF para archivos cargados manualmente ---
            # Regla: si 'name' está vacío y 'source_url' NO parece una URL (no contiene "://"),
            # usamos 'source_url' como nombre lógico del archivo (ZIP o PDF manual).
            if "name" in df_out.columns and "source_url" in df_out.columns:
                # Aseguramos que sean cadenas para aplicar condiciones
                src_series = df_out["source_url"].astype(str)
                name_series = df_out["name"].astype(str)

                mask_manual_file = (
                    src_series.notna()
                    & src_series.str.strip().ne("")         # source_url no vacío
                    & ~src_series.str.contains("://", regex=False)  # no es una URL http(s)
                    & (name_series.str.strip() == "")       # name está vacío
                )

                if mask_manual_file.any():
                    df_out.loc[mask_manual_file, "name"] = src_series[mask_manual_file].apply(
                        lambda s: Path(s).name  # sólo nombre de archivo (ej. MiZIP.zip)
                    )

            if "Fila_Excel" in df_out.columns:
                df_out = df_out.sort_values(["Fila_Excel", "Status"]).reset_index(drop=True)
            else:
                df_out = df_out.reset_index(drop=True)

            try:
                df_out["Tipo_Problema"] = df_out.apply(_infer_tipo_problema, axis=1)
            except Exception:
                df_out["Tipo_Problema"] = ""

            df_out = _standardize_status_column(df_out)

            df_export = df_out.copy()
            if "Nombre del Archivo" in df_export.columns:
                df_export = df_export.rename(columns={"Nombre del Archivo": "Archivo"})

            st.session_state.status_result_df = df_out
            st.session_state.status_export_df = df_export

            progress_bar.empty()
            status_text.markdown("✅ Validación completada.")
            st.session_state["pipeline_status_done"] = True

            if df_out is not None and not df_out.empty:
                _render_status_summary(df_out)
            else:
                st.info("No se produjeron resultados para mostrar.")

        except Exception as e:
            progress_bar.empty()
            status_text.empty()
            st.error(f"Ocurrió un error durante la validación: {e}")
    else:
        df_out = st.session_state.get("status_result_df")
        if df_out is not None and not df_out.empty:
            _render_status_summary(df_out)
        else:
            st.info("Aún no se ha ejecutado la validación de links.")

    # 4.4 Descargar Status (Excel) - dentro del mismo bloque (Paso 9)
    df_ready: Optional[pd.DataFrame] = st.session_state.get("status_export_df")
    if df_ready is None:
        df_ready = st.session_state.get("status_result_df")

    if df_ready is None or df_ready.empty:
        st.warning("Primero ejecuta la validación para generar el status.")
        ui_card_close()
        return

    file_base = Path(st.session_state.status_input_filename or "reporte_link").stem
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_xlsx = f"{file_base}_STATUS_{ts}.xlsx"

    excel_bytes = _to_excel_report(df_ready)

    colD1, colD2 = st.columns([1, 3])
    with colD1:
        st.download_button(
            "⬇️ Descargar Excel Status",
            data=excel_bytes,
            file_name=out_xlsx,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    ui_card_close()

# ======================================================
# MAIN
# ======================================================

def main():
    st.set_page_config(
        page_title=APP_TITLE,
        page_icon=APP_ICON,
        layout="wide",
        initial_sidebar_state="expanded",
    )

    apply_global_styles()
    init_session_state()

    with st.sidebar:
        render_sidebar_header()

        st.radio(
            "Módulos",
            MODULES,
            index=MODULES.index(st.session_state["module"]),
            key="module_radio",
            on_change=on_change_module,
        )
        module = st.session_state["module"]

        st.markdown("---")
        with st.expander("Recomendaciones"):
            st.markdown(
                """
                - Se recomienda en el proceso "1" descargar de forma masiva entre 400 - 500 documentos/registros como límite máximo..  
                - Para más de 500 documentos/registros, lo recomendable es:  
                    Ejecutar el app en local.
                    Dividir el Excel de Url en varios archivos (por ejemplo bloques de 500 registros) y procesarlos por partes.  
                - Esto para evitar que el contenedor de Streamlit Cloud se quede sin memoria (~1 GB de RAM) 
                """
            )

    module = st.session_state["module"]

    if module == "Home":
        page_home()
    elif module == "Report Broken Link":
        page_report_broken_unificado()
    else:
        render_hero(title=module, subtitle="Módulo no encontrado.", icon="⚠️")
        st.error("Módulo seleccionado no existe.")

if __name__ == "__main__":
    main()






















